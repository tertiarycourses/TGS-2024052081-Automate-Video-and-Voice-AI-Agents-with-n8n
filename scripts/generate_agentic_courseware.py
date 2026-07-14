from __future__ import annotations

import re
from dataclasses import dataclass
from pathlib import Path
from textwrap import dedent

ROOT = Path(__file__).resolve().parents[1]
COURSEWARE = ROOT / "courseware"
ARCHIVE = COURSEWARE / "archive"

# The WSQ house identity — used on every cover page, footer and filename.
COURSE_TITLE = "Automate Video and Voice AI Agents with n8n"
COURSE_CODE = "TGS-2024052081"
VERSION = "v4.0"
VERSION_DATE = "14 July 2026"
ORG = "Tertiary Infotech Academy Pte Ltd"
UEN = "UEN: 201200696W"
# The deck filename carries the version, and the README links to it. Derive both from one
# place, or the README goes on pointing at a deck that has moved into courseware/archive/.
PPT_STEM = f"{COURSE_TITLE.replace(' ', '-')}-{VERSION}"

# prodoc.py (installed by /wsq-setup) carries the house cover page, the version
# control record, the real Word TOC field and the Page X of Y footer. Reuse it
# rather than hand-rolling a second, subtly-different house style.
SKILLS = ROOT / ".claude" / "skills"
PRODOC = SKILLS / "wsq-slides" / "reference" / "prodoc.py"
ASSETS = SKILLS / "tertiary-learner-guide" / "assets"
ORG_LOGO = ASSETS / "tertiary-infotech-logo.png"
COURSE_LOGO = ASSETS / "n8n-course-logo.png"

VERSION_ROWS = [
    ("v1.0", "8 July 2026", "Initial release.", "Dr Alfred Ang"),
    ("v2.0", "11 July 2026", "Agentic AI loop engineering rewrite; labs 1-6.", "Dr Alfred Ang"),
    ("v3.0", "12 July 2026",
     "Aligned 100% to the runnable labs: RAG concepts, Retell + ngrok setup, "
     "knowledge base and voice cloning, Vapi FAQ agent, Wav2Lip/MuseTalk/HeyGen "
     "lip-sync comparison, interactive avatars (in-browser and LiveAvatar), and "
     "Gemini Veo 3 video generation.", "Dr Alfred Ang"),
    ("v4.0", VERSION_DATE,
     "Rebuilt as three topics over two days, aligned to the WSQ skill "
     "Artificial Intelligence Application (AER-TEM-4026-1.1): Topic 1 Chatbot "
     "(Labs 0-3), Topic 2 Voice Agent (Labs 4-5, ElevenLabs + Vapi), Topic 3 "
     "Video Agent (Labs 6-10). Added learner introductions, WSQ learning "
     "outcomes and the GitHub lab download page; matched the v4.0 WA/CS "
     "assessment set.", "Dr Alfred Ang"),
]
LABS_DIR = ROOT / "labs_local_n8n"


@dataclass(frozen=True)
class Lab:
    topic: int
    lab_no: str
    slug: str
    title: str
    minutes: int
    build: str
    concepts: list[tuple[str, str]]
    steps: list[str]
    verify: list[str]
    troubleshooting: list[tuple[str, str, str]]
    deliverable: str
    # Optional: screenshot keys (see SCREENSHOTS) rendered into the LG and the deck.
    shots: tuple[str, ...] = ()


SHOTS_DIR = COURSEWARE / "screenshots"

# Every screenshot used by the courseware. Captured from the running labs, so the
# deck and the Learner Guide show what the learner actually sees on screen.
# key -> (filename in courseware/screenshots, caption)
SCREENSHOTS: dict[str, tuple[str, str]] = {
    "n8n-workflows": ("n8n-workflow-list.png", "The course workflows imported into local n8n."),
    "lab1-canvas": ("lab1-ai-agent-ollama.png", "Lab 1 - AI Agent: chat trigger, AI Agent and the local Ollama model."),
    "lab2-canvas": ("lab2-rag-flow.png", "Lab 2 - RAG IT Support Chatbot: ingestion, embeddings and the vector store."),
    "lab2-uploader": ("lab2-brochure-uploader.png", "The Lab 2 upload page: the learner pastes their OWN n8n webhook URL."),
    "lab3-canvas": ("lab3-cx-agent-rag.png", "Lab 3 - CX Agent with RAG: the agent plus its retrieval tool."),
    "lab3-site": ("lab3-website-home.png", "Lab 3 - Cook & Bake Academy site: the customer-facing front end."),
    "lab3-settings": ("lab3-chat-webhook-settings.png", "Lab 3 - the chat widget's gear: each learner points it at their own n8n webhook."),
    "lab4-site": ("lab4-website-home.png", "Lab 4 - GG Hair Salon site: the Book by Voice call to action."),
    "lab4-settings": ("lab4-webhook-settings.png", "Lab 4 - Settings: the ElevenLabs Web Call Trigger webhook URL and optional agent ID. Nothing is hardcoded."),
    "lab5-vapi-flow": ("lab5-vapi-flow.png", "Lab 5 - the Vapi custom-LLM flow: webhook -> AI Agent (Ollama) -> OpenAI-shaped response."),
    "lab5-vapi-site": ("lab5-vapi-site.png", "Lab 5 - MediRefill: Ava, the prescription-refill voice assistant, built on the Vapi Web SDK."),
    "lab5-vapi-settings": ("lab5-vapi-settings.png", "Lab 5 - Settings: the learner's own Vapi PUBLIC key and assistant ID. The private key never touches the browser."),
    "lab6-lipsync-studio": ("lab6-lipsync-studio.png", "Lab 6 - Digital Human Studio (http://localhost:8137): face detected, script loaded, and the avatar speaking in the live preview."),
    "lab6-lipsync-engines": ("lab6-lipsync-engines.png", "Lab 6 - the controls that matter: Draft with Ollama (gemma4), the voice engine, and the LIP SYNC renderer picker."),
    "lab7-heygen-site": ("lab7-heygen-site.png", "Lab 7 - the GG News Studio: gemma4's script in the teleprompter and an honest render progress bar."),
    "lab7os-site": ("lab7os-site.png", "Lab 7-os - the open-source News Studio: 100% free and local - n8n + Ollama + TTS + ffmpeg, no cloud, no credits."),
    "lab8-browser-avatar": ("lab8-interactive-avatar.png", "Lab 8 - Aria, rendered in the browser: speech in, gemma4 reply, and the mouth drawn live. The latency is printed, not claimed."),
    "lab9-liveavatar": ("lab9-liveavatar-site.png", "Lab 9 - the LiveAvatar embed: n8n mints a short-lived session URL, so the API key never reaches the browser."),
    "lab10-veo-site": ("lab10-veo-site.png", "Lab 10 - the Veo studio: one prompt in, gemma4 writes the shot script, Veo 3.1 renders the clip."),
    "ngrok-status": ("ngrok-status.png", "http://127.0.0.1:4040 - ngrok's own status page. URL = your public address; Addr = the local n8n it forwards to."),
}


# Screenshots for labs whose Lab(...) entry does not carry its own `shots` tuple.
LAB_SHOTS: dict[str, tuple[str, ...]] = {
    "0": ("n8n-workflows",),
    "1": ("lab1-canvas",),
    "2": ("lab2-canvas", "lab2-uploader"),
    "3": ("lab3-canvas", "lab3-site", "lab3-settings"),
    "4": ("lab4-site", "lab4-settings"),
    "5": ("lab5-vapi-flow", "lab5-vapi-site", "lab5-vapi-settings"),
    "6": ("lab6-lipsync-studio", "lab6-lipsync-engines"),
    "7": ("lab7-heygen-site",),
    "7-os": ("lab7os-site",),
    "8": ("lab8-browser-avatar",),
    "9": ("lab9-liveavatar",),
    "10": ("lab10-veo-site",),
}


def lab_shots(lab: Lab) -> tuple[str, ...]:
    return lab.shots or LAB_SHOTS.get(lab.lab_no, ())


def shot_path(key: str) -> Path:
    return SHOTS_DIR / SCREENSHOTS[key][0]


def shot_md(key: str, prefix: str) -> str:
    """Markdown image + caption, or an empty string if the file is not captured yet."""
    if key not in SCREENSHOTS or not shot_path(key).exists():
        return ""
    fname, caption = SCREENSHOTS[key]
    return f"![{caption}]({prefix}{fname})\n\n*{caption}*"


TOPICS = [
    (
        "Topic 01 - Chatbot",
        "Agents, retrieval and grounded answers: your first n8n AI agent, a RAG chatbot that admits what it does not know, and a customer-facing course advisor on a real website. Labs 0-3.",
    ),
    (
        "Topic 02 - Voice Agent",
        "Two vendors, two architectures: ElevenLabs runs the model and calls your n8n tools; with Vapi your n8n workflow IS the model. The contrast is the lesson. Labs 4-5.",
    ),
    (
        "Topic 03 - Video Agent",
        "Lip-sync, avatars and text-to-video: from a script, to a mouth that moves, to a face that talks back, to video from nothing but a sentence. Labs 6-10.",
    ),
]

# ── WSQ alignment ─────────────────────────────────────────────────────────────
# TSC: Artificial Intelligence Application (AER-TEM-4026-1.1). The chatbots,
# voice agents and video avatars built in the labs ARE the AI digital human
# applications the skill describes. One LO per topic; the WA tests K1-K6 and
# the Case Study tests A1-A6, exactly as in the assessment papers (v4.0).
WSQ_TSC = "Artificial Intelligence Application (AER-TEM-4026-1.1)"

LEARNING_OUTCOMES = [
    ("LO1",
     "Analyze the strengths, limitations, and feasibility of AI digital human "
     "technology within industry contexts.",
     "K1, K6, A3, A5"),
    ("LO2",
     "Evaluate the performance of AI digital human applications and analyze "
     "their effectiveness.",
     "K2, K3, A1, A4"),
    ("LO3",
     "Assess the design and improvements for AI digital human technology.",
     "K4, K5, A2, A6"),
]

TOPIC_LO = {1: LEARNING_OUTCOMES[0], 2: LEARNING_OUTCOMES[1], 3: LEARNING_OUTCOMES[2]}

KNOWLEDGE_STATEMENTS = [
    ("K1", "Range of AI applications"),
    ("K2", "Concepts pertaining to performance effectiveness and analysis"),
    ("K3", "Methods of evaluating effectiveness of AI applications"),
    ("K4", "Algorithm design and implementation"),
    ("K5", "Methods of evaluating process improvements to the engineering processes using AI"),
    ("K6", "Applicability of AI in the industry"),
]

ABILITY_STATEMENTS = [
    ("A1", "Analyse algorithms in the AI applications"),
    ("A2", "Establish the correlation between design of algorithms and efficiency"),
    ("A3", "Identify strengths and limitations of the AI applications"),
    ("A4", "Evaluate various AI applications to compare strengths and limitations of the AI applications"),
    ("A5", "Assess feasibility of AI applications to the engineering processes"),
    ("A6", "Assess improvements on the engineering and maintenance processes"),
]

GITHUB_URL = "https://github.com/tertiarycourses/TGS-2024052081-Automate-Video-and-Voice-AI-Agents-with-n8n"


LABS = [
    # ── Topic 1 - Chatbot (Day 1 morning) ────────────────────────────────────
    Lab(
        1,
        "0",
        "setup-local-n8n",
        "Set Up n8n Locally",
        45,
        "A working local stack: Docker running n8n and Postgres, Ollama serving gemma4 and nomic-embed-text, all talking to each other.",
        [
            ("Docker Compose", "One file starts n8n and Postgres as a repeatable local stack."),
            ("Ollama", "Runs the chat model (gemma4) and the embedding model on your machine - free, offline, private."),
            ("host.docker.internal", "How n8n inside a container reaches Ollama on the host. Inside the container, localhost means the container itself."),
            ("The Active rule", "A production /webhook/ path exists only when a workflow is switched Active. Inactive means 404, every time."),
        ],
        [
            "Install Docker Desktop and Ollama (macOS: `brew`, Windows: `winget`), then launch Docker and wait for the whale to settle.",
            "From `labs_local_n8n/lab0`, run `docker compose pull`, then `docker compose up -d`, and confirm both containers with `docker compose ps`.",
            "Pull the two models: `ollama pull gemma4` and `ollama pull nomic-embed-text`, then confirm both appear in `ollama list`.",
            "Open `http://localhost:5678` and create the n8n owner account. It is local only, and there is no password reset - write it down.",
            "Create an Ollama credential named `Ollama local` with base URL `http://host.docker.internal:11434`, and press Test.",
            "Learn the everyday commands: `docker compose stop` to pause, `up -d` to resume - and never `down -v` unless you mean to erase every workflow you own.",
        ],
        [
            "`docker compose ps` shows the n8n and postgres containers running.",
            "n8n opens at `http://localhost:5678` and `ollama list` shows both models.",
            "The `Ollama local` credential test passes.",
            "You can explain why the credential must not use `localhost:11434`.",
        ],
        [
            ("n8n cannot connect to Ollama", "The credential uses localhost from inside Docker.", "Use `http://host.docker.internal:11434`."),
            ("Command not recognized (Windows)", "PowerShell was already open during the install.", "Close and reopen PowerShell - a new program is not on your PATH until you do."),
            ("Port 5678 is busy", "An older n8n container is still running.", "Run `docker ps`, stop the old container, or change the compose port."),
        ],
        "A running local stack: screenshots of docker compose ps, ollama list and the passing credential test.",
    ),
    Lab(
        1,
        "1",
        "first-ai-agent",
        "Your First AI Agent",
        45,
        "The smallest possible working agent: chat trigger, AI Agent node, and local gemma4 - with a system prompt you control.",
        [
            ("Chat trigger", "n8n's built-in chat window - a test surface with no website needed."),
            ("AI Agent node", "Holds the system prompt, the memory and, later, the tools."),
            ("System prompt", "The standing instruction that shapes every reply - your main control surface."),
            ("Execution trace", "What the model actually received and returned. A fluent answer is not evidence; the trace is."),
        ],
        [
            "Import `lab1/ai-agent.json`. Open the Ollama Chat Model node and confirm the credential is `Ollama local` and the model is `gemma4:latest`.",
            "Click Chat at the bottom of the canvas and say hello. Confirm a reply arrives.",
            "Open the execution and read exactly what the model received - the whole prompt, not just your message.",
            "Open the AI Agent node and set a system prompt: 'You are a terse assistant. Never use more than two sentences.'",
            "Ask the same question again and watch the behaviour change.",
            "Change one more instruction and re-test. One variable at a time - that habit is the course.",
        ],
        [
            "The agent replies in the chat panel.",
            "The execution list shows the run, and you can read what the model received.",
            "Changing the system prompt visibly changes the behaviour.",
        ],
        [
            ("Empty reply", "The model name does not match what ollama list shows.", "Select the exact tag, e.g. `gemma4:latest`."),
            ("Connection refused", "Ollama is not running, or the base URL is wrong.", "Start Ollama and use `host.docker.internal:11434` in the credential."),
            ("The first reply is very slow", "A 9-GB model is cold-loading into memory.", "Expected once per session - later replies are fast."),
        ],
        "A working local agent plus a one-line note on how the system prompt changed its behaviour.",
    ),
    Lab(
        1,
        "2",
        "rag-it-support-chatbot",
        "RAG IT Support Chatbot",
        60,
        "A chatbot that answers only from an uploaded IT FAQ PDF - and admits it when the answer is not in the document.",
        [
            ("Ingest", "Split the PDF into chunks, embed each chunk with nomic-embed-text, store the vectors."),
            ("Retrieve", "Embed the question the same way and fetch the nearest chunks - search by meaning, not keywords."),
            ("Grounded generation", "gemma4 answers ONLY from the retrieved chunks, or says it does not know."),
            ("The refusal", "A question outside the PDF must get an honest refusal, not an invention. This is the graded behaviour."),
        ],
        [
            "Import `lab2/rag-flow.json` and set it Active - both webhooks (`/rag-upload`, `/rag-chat`) exist only while it is Active.",
            "Serve the page: from `labs_local_n8n/lab2`, run `python3 -m http.server 8092` (Windows: `python`).",
            "Open `http://localhost:8092` - never by double-clicking index.html - and upload `it-faq.pdf`. Wait for the confirmation.",
            "Ask a question the FAQ answers, and check the reply against the PDF.",
            "Ask a question it does NOT answer - 'What is the CEO's salary?' - and demand a refusal.",
            "Open the execution and read the retrieved chunks for both questions. If the right chunk never came back, no prompt rewrite will fix the answer.",
        ],
        [
            "A question covered by the PDF gets a grounded, correct answer.",
            "A question not covered gets a refusal, not an invention.",
            "You can point to the retrieved chunks in the execution trace.",
        ],
        [
            ("Every button dies silently", "The page was opened from file://.", "Serve it over http://localhost - the browser blocks fetch() on file URLs."),
            ("The website cannot reach n8n", "The workflow is not Active.", "Activate it - the production /webhook/ path does not exist until you do."),
            ("Answers are wrong or missing", "The right chunk was never retrieved.", "Read the retrieved chunks first; fix chunking or top-k, not the prompt."),
        ],
        "A RAG chatbot answering from the IT FAQ, plus one provoked refusal captured in the execution trace.",
    ),
    Lab(
        1,
        "3",
        "cx-agent-cook-and-bake",
        "CX Agent with RAG (Cook & Bake Academy)",
        60,
        "A customer-facing course advisor: a chat widget on a real-looking school website, grounded in the school's own course brochures.",
        [
            ("Brochure knowledge base", "The vector store is filled from real course brochures - the agent speaks for a business, not a toy FAQ."),
            ("Customer tone", "A CX agent answers as the business: specific courses, prices, schedules and next steps."),
            ("Conversation memory", "A follow-up like 'how much is that one?' must resolve against the previous turn."),
            ("Honest limits", "Asked about a course that does not exist, the agent says so. No invented courses, no invented prices."),
        ],
        [
            "Import `lab3/CX Agent with RAG.json` and set it Active.",
            "Open `lab3/upload-brochures.html` and upload the PDFs from `lab3/brochures/` - this fills the vector store through `/brochure-upload`.",
            "Serve the site: from `labs_local_n8n/lab3/website`, run `python3 -m http.server 8093`, and open the chat widget.",
            "Ask about a course, its price and its schedule - and check every answer against a brochure.",
            "Ask a follow-up that needs memory: 'how much is that one?'",
            "Ask for a course that does not exist, and confirm the agent admits it rather than inventing one.",
        ],
        [
            "The agent answers from the brochures, with specifics.",
            "A non-existent course gets an honest 'we don't offer that'.",
            "The chat has memory - the follow-up resolves correctly.",
        ],
        [
            ("The widget cannot reach n8n", "The workflow is inactive, or the webhook URL was changed.", "Activate the flow - the widget already points at /webhook/cx-agent."),
            ("Vague answers, no prices", "The brochures were never uploaded.", "Run upload-brochures.html first and watch the upload executions succeed."),
            ("Follow-ups lose the thread", "Memory is not wired into the agent.", "Check the memory node and its session key in the workflow."),
        ],
        "A working course-advisory chatbot on the site, with one grounded answer and one honest refusal in the trace.",
    ),
    # ── Topic 2 - Voice Agent (Day 1 afternoon) ──────────────────────────────
    Lab(
        2,
        "4",
        "voice-booking-elevenlabs",
        "Voice Booking Agent with ElevenLabs (GG Hair Salon)",
        120,
        "Nina, a voice receptionist who answers by voice, checks a real Google Calendar for free slots, and books a real appointment.",
        [
            ("Signed URL", "n8n asks ElevenLabs for a short-lived signed URL server-side, using the xi-api-key - the key never reaches the browser."),
            ("Agent tools", "check_availability and book_appointment are n8n webhooks that ElevenLabs' SERVERS call during the conversation."),
            ("Call direction", "Browser-to-n8n needs no tunnel; ElevenLabs-to-n8n needs ngrok. Before debugging any webhook, ask: who is dialling?"),
            ("Grounded voice", "The salon handbook PDF in the agent's Knowledge Base is why Nina quotes real prices instead of inventing them."),
        ],
        [
            "Import BOTH flows - `elevenlabs-web-call-flow.json` and `elevenlabs-booking-tools-flow.json` - and set both Active.",
            "On the Get Signed URL node, add a Header Auth credential named `ElevenLabs API`: name `xi-api-key`, value your key.",
            "Connect YOUR OWN Google account on the two Google Calendar nodes - the flow ships without a calendar credential on purpose.",
            "In the ElevenLabs dashboard, create the agent (Nina), upload `knowledge-base/gg-hair-salon-handbook.pdf` to her Knowledge Base, and put her agent ID into the page's Settings.",
            "Start the tunnel with `ngrok http 5678`, and register the two tools in the agent as `https://<id>.ngrok-free.app/webhook/check-availability` and `.../book-appointment`.",
            "Serve the site (`python3 -m http.server 8090` from `lab4/website`), click Book by Voice, and book a Thursday 2 PM haircut end to end.",
        ],
        [
            "Nina answers by voice and quotes handbook prices.",
            "Asked for a taken slot, she offers an alternative.",
            "The booking appears in YOUR Google Calendar - that event is the evidence.",
            "A tool execution appears in n8n DURING the call.",
        ],
        [
            ("Nina says 'let me check that' and stalls forever", "Her tool call went to localhost, which ElevenLabs' servers cannot see.", "Register the tools with your ngrok URL, and keep the tunnel window open all session."),
            ("The call never starts", "The web-call flow is inactive, or the Settings URL is wrong.", "Activate the flow and paste the PRODUCTION URL - /webhook/, never /webhook-test/."),
            ("She invents prices", "The handbook is not attached to the agent.", "Upload the PDF to the agent's Knowledge Base and save the agent again."),
        ],
        "A real calendar booking made by voice: the calendar event, the call transcript, and the n8n tool execution.",
    ),
    Lab(
        2,
        "5",
        "vapi-faq-medirefill",
        "Grounded FAQ Voice Agent with Vapi (MediRefill)",
        90,
        "Ava, a pharmacy refill assistant whose brain is YOUR n8n workflow - grounded in six FAQ topics and hard-refusing all medical advice.",
        [
            ("Custom LLM", "Vapi does speech-to-text and the voice, then calls your n8n webhook as its model. Here, n8n IS the brain."),
            ("Public vs private key", "The page gets the Vapi PUBLIC key only - it can merely start calls. The private key manages your account and never leaves the dashboard."),
            ("Grounding", "Ava answers only from the six FAQ topics in her prompt - delivery, refills, collection, payment."),
            ("The safety boundary", "Dose, interaction, substitution, symptom: one fixed refusal sentence and a pharmacist callback. Fixed wording in the prompt - never the model's judgement."),
        ],
        [
            "Import `lab5/vapi-faq-flow.json` and set it Active. Read `ava-assistant-prompt.md` - the guardrail lives in its wording.",
            "Start the tunnel: `ngrok http 5678`. Vapi's servers must reach your n8n; localhost is invisible to them.",
            "Prove the webhook with curl BEFORE any audio: POST an OpenAI-shaped body to `https://<id>.ngrok-free.app/webhook/vapi-faq` and read a real answer back.",
            "In Vapi, create an assistant whose model is that Custom LLM URL.",
            "Serve the site (`python3 -m http.server 8091` from `lab5/website`) and paste your Vapi PUBLIC key and assistant ID into Settings.",
            "Run the graded calls: 'When will my refill arrive?', then 'Can I take two instead of one?', then 'I'm having chest pains.'",
        ],
        [
            "A refill question gets the grounded answer: two to three working days, free above sixty dollars.",
            "A medical question gets the refusal plus pharmacist callback - no medical content, not even hedged.",
            "An emergency gets the escalation: 995 / A&E.",
            "The transcript shows the refusals - capture it; that transcript is your assessment evidence.",
        ],
        [
            ("Ava is silent mid-call", "Vapi cannot reach the Custom LLM URL.", "Open ngrok's Inspect tab (127.0.0.1:4040) - if nothing arrived, fix the URL or the tunnel."),
            ("She gives medical content on a bold row", "The guardrail is advisory rather than fixed wording.", "Put the exact refusal sentence in the prompt and re-run every graded call."),
            ("The page rejects the key", "The PRIVATE key was pasted into Settings.", "Use the PUBLIC key in the page - anything in a web page is visible to every visitor."),
        ],
        "A live Vapi call transcript showing grounded answers, a hard medical refusal, and the emergency escalation.",
    ),
    # ── Topic 3 - Video Agent (Day 2) ────────────────────────────────────────
    Lab(
        3,
        "6",
        "lipsync-face-off",
        "Lip-Sync Face-Off: MuseTalk vs HeyGen",
        60,
        "The same gemma4 script and the same portrait rendered by MuseTalk (local) and HeyGen (cloud), judged side by side with your own eyes.",
        [
            ("Lip-sync rendering", "An engine animates a still portrait to match speech audio - the mouth is generated, the photo is real."),
            ("Local vs cloud", "MuseTalk: free, private, ~75 s, mouth only. HeyGen: credits, face uploaded, ~40 s, the head moves and blinks."),
            ("Served, never file://", "The studio runs at http://localhost:8137 via start.command / start.bat. Double-click index.html and fetch() silently dies."),
            ("Judging with evidence", "The deliverable is a scored comparison of what you watched - not the vendor's marketing page."),
        ],
        [
            "From `labs_local_n8n/lab6`, run `./setup.sh` (app + HeyGen), or `./setup.sh --musetalk` to also download the ~3.5 GB MuseTalk weights - once.",
            "Launch with `start.command` (macOS) or `start.bat` (Windows) - it serves the Digital Human Studio at `http://localhost:8137`.",
            "Let gemma4 draft the news script in the studio - Draft with Ollama.",
            "Render the script and portrait through MuseTalk, and time it.",
            "Render the SAME script and portrait through HeyGen (API key + credits), and time it.",
            "Score both clips: mouth realism, head movement, speed, privacy, cost - and write one sentence on when you would choose each.",
        ],
        [
            "Both clips render from the identical script and portrait.",
            "You can state, in one sentence each, when you would choose each engine - backed by what you saw.",
        ],
        [
            ("'Could not load that sample'", "The studio was opened from file://.", "Always launch via start.command / start.bat - the browser blocks fetch() on file URLs."),
            ("MuseTalk render fails", "Weights not downloaded, or no GPU/MPS available.", "Re-run ./setup.sh --musetalk and read the service log."),
            ("HeyGen render rejected", "Missing API key or exhausted credits.", "Add the key in the studio settings and check the credit balance."),
        ],
        "Two rendered clips of one script, a scored comparison, and a one-line recommendation per engine.",
    ),
    Lab(
        3,
        "7",
        "avatar-news-heygen",
        "Avatar News Video with HeyGen (GG News Studio)",
        45,
        "Facts in, broadcast out: gemma4 writes spoken copy, HeyGen renders a presenter reading it, and the page polls until the video plays.",
        [
            ("Two texts", "You type FACTS; gemma4 writes SPOKEN COPY. The renderer speaks every character it is given - feed it raw facts and the avatar reads out a list."),
            ("Generate + poll", "Rendering is far too slow for one HTTP request: /heygen-generate returns a video ID at once, and the page polls /heygen-status."),
            ("Spoken-copy rules", "No bullets, no URLs, no markdown - a stray asterisk becomes an audible 'asterisk'."),
            ("Honest progress", "HeyGen reports only processing or completed, so the bar estimates elapsed time and only reaches 100% when the render truly finishes."),
        ],
        [
            "Import `lab7/heygen-news-avatar-flow.json` and set it Active.",
            "Add the `HeyGen API` credential, and confirm the `Ollama local` credential on the script node.",
            "Serve the site from `lab7/website` and open the GG News Studio.",
            "Write YOUR OWN news item - a story from your industry, your company or today's headlines - and generate; watch gemma4's spoken script appear in the teleprompter.",
            "Watch the page poll `/heygen-status` until the video is completed, then play it. Read the script aloud: does it SOUND spoken?",
            "SHARE IT: play your broadcast to the class, and note one piece of feedback you would act on.",
        ],
        [
            "The script reads as speech - no bullets, URLs or markdown.",
            "The page polls and eventually plays the finished video.",
            "YOUR OWN news broadcast is generated and played to the class.",
        ],
        [
            ("The request times out", "The flow tried to render inside a single HTTP request.", "Keep the split shape: generate returns the ID immediately; the page polls status."),
            ("The avatar says 'asterisk'", "Markdown leaked into the script.", "Tighten the system prompt: spoken copy only, no formatting characters."),
            ("Status never reaches completed", "Wrong video ID, or HeyGen credits ran out.", "Read the status execution in n8n and check the HeyGen account."),
        ],
        "Your own news broadcast, played to the class, plus quality-review notes on the script and the render.",
    ),
    Lab(
        3,
        "7-os",
        "avatar-news-open-source",
        "Open-Source News Avatar - Free and Local",
        45,
        "The same news video with zero cloud and zero credits: TTS, Wav2Lip and ffmpeg on your own machine, driven by n8n.",
        [
            ("Local render service", "start.command / start.bat runs the pipeline - TTS to speech.wav, Wav2Lip lip-sync, ffmpeg to 1920x1080 - at localhost:8099/render."),
            ("n8n calls OUT", "n8n reaches your machine at host.docker.internal - the opposite direction from Lab 4's tools, so no tunnel is needed."),
            ("One-response shape", "A ~16 s local render is fast enough to answer in a single response - no video ID, no polling."),
            ("The trade", "Wav2Lip's mouth is generated at 96x96 and soft at 1080p - but the render is free, private and fast."),
        ],
        [
            "Start the render service FIRST: `start.command` (macOS) or `start.bat` (Windows), and leave the window open.",
            "Import `lab7-opensource/os-news-avatar-flow.json` and set it Active.",
            "Open the lab website and submit the same facts you used in Lab 7.",
            "Watch the execution: n8n calls the render service at `host.docker.internal:8099/render` and waits for the finished file.",
            "Play the finished MP4 in the page - produced entirely on your machine.",
            "Compare it against the HeyGen clip from Lab 7: quality, speed, privacy, cost.",
        ],
        [
            "A finished MP4 plays in the page, rendered locally.",
            "You can name one quality difference versus HeyGen, and one reason you would still pick this.",
        ],
        [
            ("n8n cannot reach the render service", "It called localhost instead of the host machine.", "Use host.docker.internal:8099 - n8n runs inside Docker."),
            ("Nothing renders", "The render service was never started.", "Run start.command / start.bat first and keep that window open."),
            ("Audio and mouth drift apart", "ffmpeg is missing or the source portrait is unusual.", "Install ffmpeg and start from the provided portrait."),
        ],
        "A locally rendered news video and a written HeyGen-versus-local comparison.",
    ),
    Lab(
        3,
        "8",
        "interactive-avatar-brain",
        "Interactive Avatar Brain (Aria, In-Browser)",
        45,
        "A talking avatar you can interrupt: speech in, an AI Agent reply from local gemma4, and the mouth drawn live in the browser.",
        [
            ("AI Agent + Ollama", "The same two nodes as Lab 1 - AI Agent plus local gemma4 - now behind a webhook instead of a chat window."),
            ("Stateless by design", "The browser sends the transcript with every turn, so the flow needs no memory node - the system prompt carries the conversation."),
            ("Spoken register", "One or two spoken sentences, no markdown, no URLs - Aria's prompt enforces it, and a cleanup node strips anything that slips through."),
            ("Latency is the feature", "An interactive avatar lives or dies on response time. The page prints the measured reply time - watch it, do not guess it."),
        ],
        [
            "Import `lab8/avatar-chat-flow.json` and set it Active. The shape: Webhook -> Build Prompt -> AI Agent (Ollama gemma4) -> Make It Speakable -> Respond.",
            "Open the AI Agent node and read Aria's system prompt: the voice rules, the conversation rules, and the facts she may state.",
            "Open Make It Speakable and see what it strips before the avatar speaks: markdown, URLs, and a thinking model's private reasoning.",
            "Serve the `lab8` website and talk to Aria.",
            "Watch the timing panel the page prints after each reply, and screenshot it - that is your evidence.",
            "Tighten one prompt rule (for example: always end with one short follow-up question), re-test, and note what changed.",
        ],
        [
            "Aria answers helpfully by voice - never 'I did not catch that' when the message arrived.",
            "Replies are one or two spoken sentences - no markdown, no URLs.",
            "The latency-panel screenshot is captured.",
        ],
        [
            ("Aria gives an empty or 'hiccup' reply", "The model returned nothing - often a thinking model spending its budget on private reasoning.", "Open the execution and read the AI Agent output; the cleanup node's think-stripping must stay in place."),
            ("The first reply takes forever", "The local model is cold-loading into memory.", "Expected once per session - later replies are much faster. Close heavy apps if it persists."),
            ("Replies read like essays", "The system prompt's voice rules were weakened.", "Restore the rules: one or two spoken sentences, no markdown, no URLs."),
        ],
        "A conversation with Aria, the latency-panel screenshot, and one prompt improvement you made and re-tested.",
    ),
    Lab(
        3,
        "9",
        "interactive-avatar-liveavatar",
        "Interactive Avatar Session (Nova, HeyGen LiveAvatar)",
        45,
        "A cloud interactive avatar embedded in the page through a short-lived session URL minted by n8n - the API key never reaches the browser.",
        [
            ("Session minting", "n8n calls LiveAvatar with the API key server-side and returns only a short-lived embed URL to the page."),
            ("Cloud quality", "Nova looks far better than Aria - photoreal, fluid - and costs credits for every session."),
            ("The comparison", "Aria: free, private, ~2 s, drawn mouth. Nova: photoreal, cloud, credits. Both are defensible - be able to say when."),
            ("The same key instinct", "Lab 4's signed URL, Lab 9's embed URL, Lab 10's proxy URL - the browser only ever receives short-lived, single-purpose tokens."),
        ],
        [
            "Import `lab9/liveavatar-session-flow.json` and set it Active.",
            "Add the `LiveAvatar API` credential in n8n.",
            "Serve the `lab9` website and start a session with Nova.",
            "Hold a short conversation and note the quality and the latency.",
            "Write down one thing Nova does better than Aria, and one thing Aria does better than Nova.",
            "Open the execution and find the session URL - confirm no API key ever reached the page.",
        ],
        [
            "Nova loads and holds a conversation.",
            "The Aria-versus-Nova scorecard has one honest entry in each column.",
        ],
        [
            ("The session fails to start", "The LiveAvatar credential is missing or invalid.", "Fix the credential in n8n and read the execution error."),
            ("The embed loads, then dies", "The short-lived session URL expired.", "Mint a fresh session - expiring quickly is the point of these URLs."),
            ("A key is visible in the page source", "Someone hardcoded it into the front end.", "Remove it - the page must only ever receive the embed URL."),
        ],
        "A working Nova session plus the two-column scorecard comparing her against the browser-rendered avatar.",
    ),
    Lab(
        3,
        "10",
        "veo3-video-generation",
        "AI Video Generation with Gemini Veo 3 (Veo Studio)",
        60,
        "One sentence in, an 8-second cinematic clip with sound out: gemma4 writes the shot prompt, Veo 3.1 renders it, n8n proxies the file.",
        [
            ("Shot prompt", "gemma4 turns your idea into camera, lighting and motion language - a shot, not a summary."),
            ("Generate, poll, fetch", "The same long-running shape as Lab 7 - plus a third step to fetch the finished file."),
            ("The /veo-file proxy", "Google's download URL demands your API key, so n8n hands the page a URL pointing back at ITSELF, fetches the MP4 server-side, and streams it through."),
            ("Prompt iteration", "Change the shot wording, regenerate, and record what changed - that note is the deliverable."),
        ],
        [
            "Import `lab10/veo3-video-flow.json` and set it Active.",
            "Add the `Gemini API` credential, and confirm the `Ollama local` credential.",
            "Serve the `lab10` website and open the Veo Studio.",
            "Type YOUR OWN idea - one sentence, any subject you like - and generate; then read the shot prompt gemma4 actually wrote.",
            "Wait for the poll to complete and play the clip IN the page, through the /veo-file proxy. Refine the prompt once and regenerate.",
            "SHARE IT: present your best clip to the class, and say what changed between your two prompt versions.",
        ],
        [
            "The shot prompt names camera, lighting and motion - not a summary.",
            "The clip plays in the page.",
            "YOUR OWN clip is presented to the class, with the before/after prompt note.",
            "You can explain why the page never sees the Gemini key.",
        ],
        [
            ("The clip will not play in the page", "The page tried Google's URL directly and was refused.", "Play it through the /veo-file proxy - n8n adds the key server-side."),
            ("Generation fails immediately", "API quota, or a malformed prompt payload.", "Run a short test prompt and read the provider response in the execution."),
            ("The clip ignores the idea", "The shot prompt drifted from the subject.", "Tighten the subject, camera and lighting wording, then regenerate."),
        ],
        "Your own Veo clip presented to the class, the shot prompt that produced it, and a note on what changed between two prompt versions.",
    ),
]


def rag_concepts_section() -> str:
    """Topic 02 prose: what RAG, tokenization and embeddings actually are."""
    return normalize_md(dedent(
        """
        ### The concepts behind RAG

        Read this before Lab 2. The labs will work if you only click the nodes, but you cannot *debug* a RAG agent - or explain to a manager why it answered wrongly - without these four ideas.

        #### What RAG is, and the problem it solves

        A language model only knows what was in its training data. It has never seen your IT FAQ, your course brochures, or your salon's cancellation policy. Ask it anyway and it will not say "I don't know" - it will produce fluent, confident, invented text. That failure has a name: **hallucination**.

        **Retrieval-Augmented Generation (RAG)** fixes this by changing the question you ask the model. Instead of:

        > "What is the refund policy?"

        the workflow silently asks:

        > "Here are three passages from the company handbook. Using ONLY these passages, answer: what is the refund policy? If the passages do not contain the answer, say you do not know."

        The model stops being a source of facts and becomes a *reader* of facts you supply. That is the whole idea. Everything else - embeddings, vector stores, chunking - exists only to answer one narrow question: **which passages should we paste in front of the question?**

        A RAG system therefore has two phases:

        | Phase | When it runs | What happens |
        |---|---|---|
        | **Indexing** (write) | Once, when a document is uploaded | Split the document into chunks -> embed each chunk -> store the vectors |
        | **Retrieval** (read) | On every question | Embed the question -> find the nearest chunks -> paste them into the prompt -> generate |

        In your n8n workflow these are the two paths you can literally see on the canvas: the upload path that ends at the vector store, and the chat path that reads from it.

        #### Tokenization - how text becomes numbers

        Models do not read characters or words. They read **tokens**: the sub-word units the model's vocabulary is built from. A tokenizer splits text deterministically:

        ```text
        "The salon is closed on Sundays."
          -> ["The", " salon", " is", " closed", " on", " Sund", "ays", "."]
             8 tokens
        ```

        Useful rules of thumb for English: **1 token is roughly 4 characters, or about 0.75 of a word** - so 1,000 tokens is roughly 750 words, or about 1.5 pages. Rare words, names, code and non-English text split into more tokens than you would expect (`Sundays` above became two).

        Tokens matter for three practical reasons:

        1. **Context windows are measured in tokens.** Everything you paste in - the system prompt, the retrieved chunks, the conversation memory, the question - competes for the same budget. Retrieve too many chunks and you push out the instructions.
        2. **Cost and latency are measured in tokens.** Doubling the retrieved text roughly doubles the prompt cost of every single call.
        3. **Chunk size is measured in tokens** (or characters, as an approximation of them). This is the number the RAG flows in Labs 2 and 3 are built around.

        #### Embeddings - meaning as coordinates

        An **embedding** is a list of numbers (a **vector**) that represents the *meaning* of a piece of text. An embedding model reads the text and outputs a fixed-length vector - in this course, `nomic-embed-text` running in Ollama, which outputs **768 numbers** for any input, whether it is three words or three paragraphs:

        ```text
        "How do I reset my password?"  ->  [0.021, -0.118, 0.334, ... ]   768 numbers
        ```

        The magic property is that **texts with similar meaning land close together in that 768-dimensional space, even when they share no words at all.** "How do I reset my password?" sits near "I forgot my login credentials" and far from "What are your opening hours?" - which is exactly what keyword search cannot do.

        Closeness is measured with **cosine similarity**: the cosine of the angle between two vectors, from `1.0` (identical direction/meaning) through `0.0` (unrelated) to `-1.0` (opposite). Retrieval is then embarrassingly simple:

        1. Embed the user's question with the **same** model used for the documents.
        2. Compare that vector against every stored chunk vector.
        3. Return the **top-k** most similar chunks (k is typically 3 to 5).

        Two consequences follow directly, and both cause real bugs:

        - **You must embed questions and documents with the same model.** Vectors from different models are not comparable - the numbers mean different things. Switch the embedding model and you must re-index every document.
        - **A vector store is not a database you can query with SQL.** It answers only one kind of question: "what is near this vector?"

        #### Chunking - why documents are cut up

        You cannot embed a 40-page PDF as one vector. A single vector would average away all the detail, and the retrieved passage would be far too big to paste into a prompt. So the document is **split into chunks** (say 800 characters each) with a small **overlap** (say 100 characters) carried between neighbours so a sentence cut in half still appears whole in one of them.

        Chunk size is a genuine trade-off, and it is the main thing you will tune:

        | Chunks too small | Chunks too large |
        |---|---|
        | Retrieved passage lacks context; the model sees half an answer | Retrieved passage contains the answer plus three irrelevant sections |
        | High precision, low recall | High recall, low precision |
        | Model says "the document does not say" when it does | Model gets distracted and cites the wrong part; tokens are wasted |

        There is no universally correct value. That is exactly why Labs 2 and 3 make you read the retrieved chunks in the execution trace instead of guessing.

        #### Putting it together

        ```text
        INDEXING   PDF -> split into chunks -> embed each chunk -> store 768-dim vectors
                                                (nomic-embed-text)      (Supabase / Qdrant / Pinecone)

        RETRIEVAL  question -> embed question -> cosine-similarity search -> top-k chunks
                                                                              |
                                                 "Using ONLY this context: <chunks>  Q: <question>"
                                                                              |
                                                                        gemma (Ollama) -> grounded answer
        ```

        **The habit to build:** when a RAG agent answers badly, do not immediately rewrite the prompt. First look at *what was retrieved*. Open the n8n execution and read the chunks that came back from the vector store. If the right chunk was never retrieved, no prompt in the world will save the answer - the bug is in chunking, embedding, or top-k, not in the wording.
        """
    ))


def ngrok_section(prefix: str) -> str:
    """Topic 04 prose: expose local n8n to the internet so a voice platform can call it."""
    return normalize_md(dedent(
        f"""
        ### Exposing n8n with ngrok - and why you must

        Labs 4 and 5 are the first time something **outside your machine** needs to call **into** n8n. Up to now every request came from your own browser, so `http://localhost:5678` worked. It will not work now.

        When Nina says *"let me check that for you"*, the request is made by **ElevenLabs' servers**, sitting in a datacenter. To them, `localhost` means *their own machine*, not yours. The request never leaves their building. The caller hears the agent stall in silence, and nothing at all appears in your n8n executions list. Same for Vapi's Custom LLM in Lab 5.

        A tunnel fixes this. It gives your local n8n a temporary public address, and forwards anything sent there to port 5678 on your laptop.

        | Who is calling n8n | Example | Needs a tunnel? |
        |---|---|---|
        | Your own browser | Labs 1-3, the Book by Voice button, `curl` | **No.** `localhost` is correct. |
        | ElevenLabs' servers | `check_availability`, `book_appointment` | **Yes.** |
        | Vapi's servers | the Custom LLM endpoint | **Yes.** |

        **Step-by-step (macOS and Windows)**

        1. Install ngrok:

            | macOS (Terminal) | Windows (PowerShell) |
            |---|---|
            | `brew install ngrok` | `winget install ngrok.ngrok` |

            On Windows, **close and reopen PowerShell** afterwards, or `ngrok` will not be found. (If winget is unavailable, download the ZIP from `https://ngrok.com/download`, unzip it, and run `.\\ngrok.exe` from the folder you unzipped it into.)

        2. Create a free account at `https://dashboard.ngrok.com/signup`, then copy your **authtoken** from *Getting Started -> Your Authtoken*. ngrok refuses to tunnel without one.
        3. Register the token once - it is saved to your ngrok config file, so you never repeat this. The command is the same on both platforms:

            ```bash
            ngrok config add-authtoken <YOUR_TOKEN>
            ```
        4. Start the tunnel, and **leave this window open** for the whole session. Closing it kills the tunnel:

            ```bash
            ngrok http 5678
            ```
        5. Read your public URL from the `Forwarding` line ngrok prints in that terminal:

            ```text
            Forwarding   https://3af5-175-156-143-249.ngrok-free.app -> http://localhost:5678
                         └────────── your public base URL ──────────┘
            ```

        6. Open ngrok's own dashboard in your browser - **bookmark this link, you will use it all day**:

            ```text
            http://127.0.0.1:4040
            ```

            It runs on your own machine (that is why the address is `127.0.0.1`, not an ngrok domain) and it has two tabs:

            - **Status** - your public URL. Look for **Tunnels -> command_line**: `URL` is the public address to paste into ElevenLabs or Vapi, and `Addr` confirms it forwards to `http://localhost:5678`. If you lose the URL, get it here - do not restart ngrok, or the URL will change.
            - **Inspect** - every request ElevenLabs or Vapi sends you, with its headers and body, as it arrives. This is the single most useful debugging tool in this topic: when the agent stalls mid-call, this page tells you instantly whether the request even reached your machine.

            An empty **Inspect** tab during a call means the request never arrived: the URL in ElevenLabs/Vapi is wrong, the tunnel is down, or the agent was not re-published after you changed it.

        {shot_md("ngrok-status", prefix)}

        **How to build the URL you paste into ElevenLabs or Vapi**

        You are gluing two halves together. n8n supplies the path; ngrok supplies the address.

        ```text
        https://3af5-175-156-143-249.ngrok-free.app  /webhook/  check-availability
        └────────── from ngrok ──────────────────┘             └── from the n8n Webhook node ──┘
        ```

        The rule: **open the Webhook node, take its Production URL, and replace `http://localhost:5678` with your ngrok address.** Everything after `/webhook/` stays exactly as n8n shows it.

        | n8n webhook node | Path | What you paste into ElevenLabs / Vapi |
        |---|---|---|
        | Webhook - check_availability | `check-availability` | `https://<your-ngrok>/webhook/check-availability` |
        | Webhook - book_appointment | `book-appointment` | `https://<your-ngrok>/webhook/book-appointment` |
        | Vapi Webhook (Lab 5) | `vapi-faq` | `https://<your-ngrok>/webhook/vapi-faq` |

        **Prove the tunnel works before you touch the voice platform**

        ```bash
        curl -X POST https://<your-ngrok>/webhook/vapi-faq \\
          -H 'Content-Type: application/json' \\
          -d '{{"model":"gpt-4o","messages":[{{"role":"user","content":"When will my refill arrive?"}}]}}'
        ```

        A real answer means the tunnel, the webhook and the agent are all healthy. If this fails, fix it here - debugging over HTTP is far easier than debugging over audio.

        **Four things that will bite you**

        | Trap | What you see | Fix |
        |---|---|---|
        | The free URL changes every restart | Calls worked this morning, fail after lunch | Re-copy the new URL into ElevenLabs/Vapi and **publish the agent again**. Start ngrok once per day and leave it running. |
        | You used the **Test** URL | Works once during your demo, never again | Use the **Production** URL - `/webhook/`, never `/webhook-test/`. |
        | The workflow is not published | 404 "not registered" | Publish/activate the workflow, then retest. |
        | You set `WEBHOOK_URL` on the n8n container | Every lab now shows learners a public URL that goes stale | Do NOT change the Docker config. `localhost` must keep working for Labs 1-3. Paste the tunnel URL only where it is needed. |

        Closing the ngrok window kills the tunnel and the public URL is gone for good - the next start hands you a different one.
        """
    ))


def video_pipeline_section() -> str:
    """Topic 05 prose: the avatar-video pipeline, and the three lip-sync engines."""
    return normalize_md(dedent(
        """
        ### How the avatar video pipeline works

        Every video lab in this topic is the same three-step pipeline. Only the *renderer* changes - and choosing the renderer is the engineering decision the topic is really about.

        ```text
        YOU                    n8n                          RENDERER            BACK TO YOU
        ---------------------------------------------------------------------------------------
        type a TOPIC  --POST-->  Webhook
        (raw facts)                |
                                   v
                          Write Script (Ollama)        gemma4 turns facts into
                          + gemma4 chat model          SPOKEN COPY (~55-70 words)
                                   |
                                   v
                          send script + portrait --->  HeyGen / Wav2Lip / MuseTalk
                                   |                   renders the talking video
                                   v
                          Respond  ---------------------------------------->  video plays
                                                                              in the page
        ```

        **The two texts are not the same text, and this is the point people miss.** You type a *topic* - "Belgium beats USA 4-1 to reach the quarter-finals". That is a fact, not something an anchor can read aloud. gemma4 turns it into *broadcast copy* - "Welcome back to Global Sports News! Massive upsets continue...". The renderer speaks **every character** it is given: feed it your raw topic and the avatar reads out a comma-spliced list of facts. The system message is what makes the difference - it forbids stage directions, markdown and camera notes, because the avatar would dutifully read them out loud.

        **Synchronous or polled?** Two shapes appear in this topic, and they change the front end:

        | Lab | Shape | Why |
        |---|---|---|
        | Lab 7 (HeyGen) | n8n returns a `video_id` at once; the **page polls** `heygen-status` | Cloud renders take 1-3 minutes. Holding an HTTP request open that long is fragile. |
        | Lab 7-os (local) | n8n returns the finished `video_url` in **one** response | A local render takes ~16 s, so the caller can simply wait. |

        Both pages show a progress bar. Note what it does NOT do: HeyGen reports only *processing* or *completed* - there is **no percentage to read**. So the bar is an honest elapsed-time estimate that eases toward 95% and stops there, reaching 100% only when the renderer actually reports done. A bar that sits at 100% while still spinning is a lie, and learners notice.

        ### The three lip-sync engines

        All three animate a still portrait to speech. They are not interchangeable, and the labs make you prove it: Lab 6 pits MuseTalk against HeyGen on the same photo and script, and Lab 7-os puts Wav2Lip to work in the free local pipeline.

        | | Wav2Lip | MuseTalk | HeyGen |
        |---|---|---|---|
        | Where it runs | your machine | your machine | the cloud |
        | Cost | free | free | **credits** |
        | Privacy | photo never leaves | photo never leaves | photo + audio uploaded |
        | Speed (measured) | **~16 s** for a 7-second clip | ~75 s (Apple M4) | ~40 s to 3 min |
        | Mouth | 96x96 model, upscaled - **soft at 1080p** | real inpainted pixels: teeth, lips, shadow | photoreal |
        | Head / blinks | **no** - the head is frozen | **no** - the head is frozen | **yes** |
        | Needs | ffmpeg + the checkpoint | a GPU (Apple MPS / NVIDIA) + 3.5 GB weights | an API key and credits |
        | Fails when | - | no GPU | no credits, no network |

        The honest summary: **Wav2Lip is the fastest, MuseTalk looks the best, and only HeyGen moves the head.** Wav2Lip's mouth is generated at 96x96 pixels and then scaled up into a 1080p frame - the *timing* is excellent (it is still the sync benchmark others are measured against), but the mouth itself is soft, and the larger your output, the more you notice. MuseTalk inpaints real mouth pixels in latent space, which is exactly why it costs about seven times the compute.

        **The workplace question is not "which is best".** It is: would you upload a customer's face to a vendor in order to gain a moving head? Would you accept a soft mouth to keep a 16-second turnaround in a classroom? There is no single right answer - there is a defensible one, and writing it down is the deliverable.
        """
    ))


def voice_setup_section(prefix: str) -> str:
    """Topic 02 prose: the two voice architectures, and where every webhook goes."""
    return normalize_md(dedent(
        f"""
        ### Two vendors, two architectures - the contrast is the lesson

        Labs 4 and 5 both put a voice on an agent, but they split the work in opposite ways. Understanding the split is what the topic teaches; everything else is form filling.

        | | **Lab 4 - ElevenLabs (GG Hair Salon)** | **Lab 5 - Vapi (MediRefill)** |
        |---|---|---|
        | Who runs the model | **ElevenLabs** | **your n8n workflow** (Vapi "Custom LLM") |
        | What n8n does | mints a signed URL, and serves the tools | **is the brain** |
        | Call path | browser -> n8n -> signed URL -> WebSocket | browser -> Vapi, then Vapi -> your n8n |
        | What the browser gets | a short-lived **signed URL** | the Vapi **public** key only |
        | What it never sees | your `xi-api-key` | your Vapi private key |

        ### Where each webhook goes, and which way it points

        Most of the pain in this topic comes from confusing webhooks that point in **opposite directions**. Before you debug anything, ask: *who is dialling?*

        | Webhook | Direction | Where the URL is written | Must it be public? |
        |---|---|---|---|
        | **Web-call trigger** (`/webhook/elevenlabs-web-call`) | Browser -> n8n -> ElevenLabs API | In the **website**: click **⚙ Settings** on the page and paste your own URL. Nothing is hardcoded. | No. `http://localhost:5678` is fine - your own browser calls it. |
        | **Agent tools** (`check_availability`, `book_appointment`) | ElevenLabs cloud -> n8n | In **ElevenLabs**, in each tool's URL field on the agent | **Yes.** ElevenLabs' servers cannot reach your `localhost`. This needs the tunnel. |
        | **Vapi Custom LLM** (`/webhook/vapi-faq`) | Vapi cloud -> n8n | In **Vapi**, as the assistant's model URL | **Yes.** Same reason. |

        ### How the ElevenLabs key stays safe

        ElevenLabs does not hand the browser a raw token. Your n8n asks ElevenLabs for a short-lived **signed URL** - server-side, using the `xi-api-key` in a Header Auth credential - and returns *only that* to the page. The browser opens a WebSocket to the signed URL and talks. **Your API key never reaches the browser.** You will meet the same instinct twice more: Lab 9's embed URL and Lab 10's proxy URL.

        {shot_md("lab4-site", prefix)}

        {shot_md("lab4-settings", prefix)}

        ### The Vapi rule that protects your account

        The MediRefill page needs two values in ⚙ Settings: your Vapi **public** key and the assistant ID. A public key can only *start calls*. Your **private** key manages your whole account - and anything pasted into a web page is visible to every visitor who opens DevTools. **Never paste the private key into the page.**

        ### Why Ava's guardrail is fixed wording, not judgement

        MediRefill is a pharmacy. A confident wrong answer in retail costs a refund; in a pharmacy it is a safety incident. So Ava's prompt (`lab5/ava-assistant-prompt.md`) hard-codes the refusal: any question about a dose, an interaction, a substitution or a symptom gets one fixed sentence and a **pharmacist callback** - and an emergency gets **995 / A&E**. If she answers a medical question with medical content - even hedged, even with a disclaimer - the guardrail failed, and the transcript that proves it is your evidence.
        """
    ))


def voice_conversation_section() -> str:
    """Topic 04 prose: the turn-by-turn call the learner runs with the voice agent."""
    return normalize_md(dedent(
        """
        ### The voice conversation, turn by turn

        Use this script for your first end-to-end call in Lab 4, once the tools and the knowledge base are wired up. Click **Book by Voice** on `http://localhost:8090`, allow the microphone, and wait for Nina to speak first. Say one line at a time and let her finish - interrupting is the most common cause of a broken slot capture.

        The **What to listen for** column is what you are grading. If a turn fails, note it, keep going, and fix one thing at a time afterwards.

        | # | You say | Nina should | What to listen for |
        |---|---|---|---|
        | 1 | *(say nothing - just listen)* | Greet and offer help: "Thanks for calling GG Hair Salon, this is Nina. How can I help you today?" | She opens the call. Silence here means the signed URL was minted but the audio session never started. |
        | 2 | "Hi Nina, I'd like to book a haircut." | Ask a single question - which service, or which day. | **One** question, not three. A multi-part question is a prompt defect. |
        | 3 | "How much is a women's cut?" | Answer **$65**, a 1-hour slot. | The price comes from the handbook. A vague or wrong price means the KB is not attached. |
        | 4 | "And what's your cancellation policy?" | State the **12-hour** rule and the **50%** charge for a late cancellation or no-show. | This fact exists only in the PDF. This is the grounding proof. |
        | 5 | "Do I need to pay a deposit for colour?" | Say **$30**, applied to the final bill. | Second grounded fact. She should not hedge. |
        | 6 | "Do you do nail extensions?" | Decline to guess and offer to check with a stylist. | The honest refusal. If she invents a nail price, the grounding instruction is missing. |
        | 7 | "Okay, book me the women's cut." | Start slot filling: ask for the day and time. | She moves from answering to acting - the KB answers questions, the tools take actions. |
        | 8 | "Thursday at 2 PM." | Check availability through the n8n tool webhook and respond. | Watch the n8n executions list: a new execution must appear **during** the call. |
        | 9 | "It's Alex, and my number is 9123 4567." | Read the details back for confirmation. | Every slot repeated: name, service, date, time, contact. |
        | 10 | "Yes, that's correct." | Confirm the booking and close the call. | The booking is created only **after** you confirm - never before. |

        **Two more calls you must run** - they are what your QA notes are built from:

        - **The incomplete caller.** At turn 8 say only *"sometime Thursday"*. Nina must ask a repair question for the time. An agent that silently invents 9:00 AM has failed.
        - **The changed-mind caller.** At turn 10 say *"actually, make it Friday instead"*. She must update the slot and re-confirm, not book Thursday anyway.

        Save all three transcripts from the conversation history in the ElevenLabs dashboard. They are the evidence for Lab 4 and for the Case Study.
        """
    ))


def lab_dir(lab: Lab) -> Path:
    return LABS_DIR / f"topic-{lab.topic:02d}" / f"lab-{lab.lab_no}-{lab.slug}"


def md_table(rows: list[tuple[str, str]]) -> str:
    body = "\n".join(f"| {a} | {b} |" for a, b in rows)
    return f"| Concept | In one line |\n|---|---|\n{body}"


def normalize_md(text: str) -> str:
    cleaned = []
    for line in text.splitlines():
        while line.startswith("    "):
            line = line[4:]
        cleaned.append(line.rstrip())
    return "\n".join(cleaned).strip() + "\n"


def lab_readme(lab: Lab) -> str:
    concepts = md_table(lab.concepts)
    steps = "\n".join(f"{i}. {step}" for i, step in enumerate(lab.steps, 1))
    verify = "\n".join(f"- [ ] {item}" for item in lab.verify)
    trouble = "\n".join(f"| {a} | {b} | {c} |" for a, b, c in lab.troubleshooting)
    agent_prompt = dedent(
        f"""
        I am completing {lab.title} in the course "Automate Video and Voice AI Agents with n8n".
        Act as a careful agentic AI workflow reviewer. Review my workflow design against this goal:
        {lab.build}

        Check for:
        - missing trigger, input, output, or credential boundary
        - weak system prompt or unclear tool description
        - missing evaluation case, refusal case, or human review gate
        - any place where a secret could leak to browser code or exported notes

        Return a concise list of fixes, then give me one improved test case.
        """
    ).strip()
    return normalize_md(dedent(
        f"""
        # Lab {lab.lab_no} - {lab.title}

        > Topic {lab.topic} - Approximately {lab.minutes} minutes - Adult workplace lab

        ## The build so far

        This lab is part of a complete agentic AI automation course. Each lab follows the same engineering loop:

        **Define -> Build -> Observe -> Evaluate -> Improve -> Guardrail -> Document**

        In this lab you build: **{lab.build}**

        ## What you will build

        {lab.build} The work is intentionally practical: by the end of the lab you should have an artifact that can be inspected, tested, and reused in the capstone.

        ## Concepts you will meet

        {concepts}

        ## Before you start

        - Complete the setup in Lab 1.1 or confirm Docker, n8n, Ollama, and your browser are ready.
        - Keep n8n executions visible while testing. The execution trace is your evidence.
        - Do not place API keys in HTML, JavaScript, Markdown examples, screenshots, or exported workflow notes.
        - Commit or export a checkpoint before making large workflow changes.

        ## Step-by-step instructions

        {steps}

        ## Agentic AI loop checklist

        1. **Define** the user, trigger, input, tool boundary, and output.
        2. **Build** the smallest workflow that proves the behavior.
        3. **Observe** the execution trace, model output, and external service response.
        4. **Evaluate** with normal, edge, and unsafe inputs.
        5. **Improve** one thing at a time and rerun the same tests.
        6. **Guardrail** credentials, refusal behavior, human handoff, and publishing actions.
        7. **Document** the final setup, test evidence, and rollback step.

        ## Vibe / agent prompt

        Paste this into your coding or workflow agent after you have a first draft:

        ```text
        {agent_prompt}
        ```

        ## Verify

        {verify}

        ## Reflection questions

        1. What did the agent or workflow do correctly on the first attempt?
        2. What evidence proves the output was grounded, safe, or complete?
        3. Which failure case was most useful for improving the workflow?
        4. What should a human still review before this automation is used with real customers?

        ## Common errors

        | Error | Likely cause | Fix |
        |---|---|---|
        {trouble}

        ## Submission evidence

        Submit the following:

        - Workflow export or screenshot of the main workflow canvas.
        - Screenshot or copy of the successful execution output.
        - One normal test case and one edge or unsafe test case.
        - A short note explaining what you changed after evaluation.

        ## Lab deliverable

        {lab.deliverable}
        """
    ))


def learner_guide(img_prefix: str = "courseware/screenshots/") -> str:
    """img_prefix makes the image links resolve from wherever the guide is written."""
    toc = [
        "# Learner Guide - Automate Video and Voice AI Agents with n8n",
        "",
        f"**Course code:** {COURSE_CODE}  |  **Conducted by:** {ORG}  |  "
        f"**Version:** {VERSION}  |  **Date:** {VERSION_DATE}",
        "",
        "## Contents",
        "",
        "- Introduction",
        "- Course learning outcomes (WSQ)",
        "- How this course uses agentic AI loop engineering",
        "- Environment setup",
        "- Topic and lab guides",
        "- How you are assessed",
        "- Preparing for the Case Study",
        "- Troubleshooting and glossary",
        "",
    ]
    intro = dedent(
        f"""
        ## Introduction

        This Learner Guide supports the adult training course **Automate Video and Voice AI Agents with n8n**. The course teaches learners how to design, build, test, and improve practical AI automations using n8n, Ollama, RAG, ElevenLabs, Vapi, HeyGen, LiveAvatar, Google Veo 3.1 and open-source lip-sync rendering - the chatbot, voice-agent and video-avatar applications collectively known as **AI digital humans**.

        The emphasis is not "click nodes until it works". The emphasis is engineering judgement. Learners will practise the agentic AI loop:

        **Define -> Build -> Observe -> Evaluate -> Improve -> Guardrail -> Document**

        Every lab produces evidence: workflow exports, screenshots, test cases, quality rubrics, generated videos, call transcripts, or monitoring runbooks. This makes the course suitable for adult learners who need workplace-ready habits, not just tool demonstrations.

        ## Course learning outcomes

        This course delivers the WSQ skill **{WSQ_TSC}**. The chatbots, voice agents and video avatars you build in the labs are the AI digital human applications the skill describes. By the end of the course, learners will be able to:

        {chr(10).join(f"- **{lo}** - {text} *({ka})*" for lo, text, ka in LEARNING_OUTCOMES)}

        Each topic delivers one learning outcome: Topic 1 (Chatbot) builds the foundation for LO1, Topic 2 (Voice Agent) for LO2, and Topic 3 (Video Agent) for LO3 - and every lab feeds evidence into all three.

        ### The knowledge and ability statements behind the outcomes

        The Written Assessment (SAQ) tests the six knowledge statements; the Case Study tests the six ability statements. Nothing is assessed that the labs did not make you do.

        | # | Knowledge (tested in the WA) | # | Ability (tested in the CS) |
        |---|---|---|---|
        {chr(10).join(f"| {k} | {kt} | {a} | {at} |" for (k, kt), (a, at) in zip(KNOWLEDGE_STATEMENTS, ABILITY_STATEMENTS))}

        ## How this course uses agentic AI loop engineering

        Agentic AI loop engineering is the repeatable practice of designing and improving an AI workflow through evidence. A model response is not trusted because it sounds fluent. A workflow is not accepted because it ran once. Each automation must define the job, expose the right tools, capture the execution trace, evaluate against test cases, improve the weakest behavior, add guardrails, and document how to operate it.

        ### The seven-part loop

        1. **Define** - Name the user, trigger, input, output, success criteria, and non-goals.
        2. **Build** - Create the smallest useful workflow before adding features.
        3. **Observe** - Inspect n8n executions, model prompts, retrieved documents, API responses, and generated media.
        4. **Evaluate** - Use normal, edge, unsafe, and unsupported inputs. Score the result.
        5. **Improve** - Change one prompt, node, chunking parameter, or guardrail at a time.
        6. **Guardrail** - Protect secrets, personal data, external publishing, bookings, refunds, and unsupported claims.
        7. **Document** - Record setup, test evidence, decisions, fallback, and owner.

        ### Adult learning approach

        The labs use workplace examples: course advisory, IT support, appointment booking, training videos, customer follow-up, and publishing. Learners are expected to compare alternatives, explain trade-offs, and build evidence. Trainers should ask learners to show execution traces and quality rubrics, not only final screens.

        ## Environment setup

        ### Required tools

        Every lab runs on **macOS and on Windows**. The tools are identical; only the way you install them and the way you type a command differ. Where this guide shows a command, it gives both.

        | Tool | What it is for |
        |---|---|
        | Docker Desktop | Runs n8n and Postgres |
        | Ollama | The local chat and embedding models - no cloud, no cost |
        | A modern browser | The lab websites; must allow microphone access for the voice labs |
        | Python 3 | Serving the lab websites over `http://localhost` |
        | ffmpeg | Local video rendering and verification (Topic 03) |
        | ngrok | Labs 4 and 5 only, when a voice platform's servers must call into your n8n |
        | Paid accounts (optional) | ElevenLabs, Vapi, HeyGen, LiveAvatar, Google Gemini (Veo 3.1) |

        ### Installing the tools

        **The package manager (do this first - everything else is one line after it)**

        | | macOS | Windows |
        |---|---|---|
        | Manager | **Homebrew** - paste the install line from `https://brew.sh` into Terminal | **winget** - already built into Windows 10/11. Nothing to install. |
        | Terminal | **Terminal** (Applications -> Utilities) | **PowerShell** - press Start, type `powershell`, and open it |

        **The tools**

        | Tool | macOS | Windows |
        |---|---|---|
        | Docker Desktop | `brew install --cask docker`, then launch Docker from Applications | `winget install Docker.DockerDesktop`, then launch Docker Desktop |
        | Ollama | `brew install ollama` then `ollama serve` | `winget install Ollama.Ollama` (it runs in the system tray) |
        | Python 3 | Already installed. Check: `python3 --version` | `winget install Python.Python.3.12` - **tick "Add python.exe to PATH"** if you use the installer instead |
        | ffmpeg | `brew install ffmpeg` | `winget install Gyan.FFmpeg` |
        | ngrok | `brew install ngrok` | `winget install ngrok.ngrok` |

        On Windows, **close and reopen PowerShell after installing** - a new program is not on your PATH until you do. If a command is "not recognized", that is almost always the reason.

        Docker Desktop on Windows needs **WSL 2**. The installer usually enables it, but if Docker refuses to start, run `wsl --install` in PowerShell **as Administrator**, reboot, and start Docker again.

        ### Core setup steps

        Both platforms, from the repository folder (get the repository itself from GitHub first: {GITHUB_URL} - Code -> Download ZIP, or `git clone`).

        **Windows: unblock the ZIP before extracting.** Downloaded files carry the *Mark-of-the-Web*, and Windows silently blocks the scripts inside. Right-click the ZIP -> **Properties** -> tick **Unblock** -> OK, *then* extract. If you already extracted and SmartScreen warns on a `start.bat`, click **More info -> Run anyway** - the launcher then unblocks the rest of that lab's files itself. (`git clone` avoids all of this.)

        Then:

        ```bash
        cd labs_local_n8n/lab0
        docker compose pull
        docker compose up -d
        ollama pull gemma4
        ollama pull nomic-embed-text
        ```

        Then open n8n at `http://localhost:5678` and create your owner account.

        In n8n, create an **Ollama** credential with this base URL:

        ```text
        http://host.docker.internal:11434
        ```

        This matters because n8n runs inside Docker. From the container's point of view, `localhost` is the container itself, not the machine running Ollama. `host.docker.internal` is how a container says "my host". It works the same on macOS and Windows.

        ### The four command differences you will actually hit

        This is the whole list. Everything else in the course is identical on both platforms.

        | Task | macOS (Terminal) | Windows (PowerShell) |
        |---|---|---|
        | Serve a lab website | `cd lab4/website`<br>`python3 -m http.server 8090` | `cd lab4\\website`<br>`python -m http.server 8090` |
        | One-click launcher | double-click `start.command` | double-click `start.bat` |
        | Path separator | `lab4/website` (forward slash) | `lab4\\website` (backslash) |
        | A `curl` with a JSON body | single quotes work:<br>`curl -X POST url -H 'Content-Type: application/json' -d '{{"a":1}}'` | PowerShell mangles quotes. Use `curl.exe` and escape:<br>`curl.exe -X POST url -H "Content-Type: application/json" -d '{{\\"a\\":1}}'`<br>Or simply use the n8n UI's own test panel instead. |

        Note the Python one: on macOS the command is **`python3`**; on Windows it is **`python`**. Typing `python3` on Windows opens the Microsoft Store, which is confusing and does nothing useful.

        ### Course evidence folder

        Create a local folder outside the repository for personal evidence:

        ```text
        course-evidence/
          lab-01/
          lab-02/
          screenshots/
          workflow-exports/
          rubrics/
          videos/
        ```

        Do not store API keys in this folder.
        """
    ).strip()

    lab_sections = []
    for idx, (topic_title, topic_desc) in enumerate(TOPICS, 1):
        lab_sections.append(f"\n\n## {topic_title}\n\n{topic_desc}\n")
        topic_labs = [lab for lab in LABS if lab.topic == idx]
        lab_sections.append("\n### Key concepts\n\n")
        lab_sections.append(
            "\n".join(
                f"- **{lab.title}:** {lab.build}" for lab in topic_labs
            )
        )
        lab_sections.append("\n\n")
        # Topic 1 needs RAG/tokenization/embeddings defined before the RAG labs.
        if idx == 1:
            lab_sections.append("\n" + rag_concepts_section() + "\n")
        # Topic 2 needs the two voice architectures and the tunnel explained
        # before either voice lab makes sense.
        if idx == 2:
            lab_sections.append("\n" + voice_setup_section(img_prefix) + "\n")
            lab_sections.append("\n" + ngrok_section(img_prefix) + "\n")
        # Topic 3 needs the pipeline shape and the engine trade-off up front.
        if idx == 3:
            lab_sections.append("\n" + video_pipeline_section() + "\n")
        for lab in topic_labs:
            lab_sections.append("\n\n")
            lab_sections.append(
                dedent(
                    f"""

                    ### Lab {lab.lab_no} - {lab.title}

                    **Time:** {lab.minutes} minutes

                    **Goal:** {lab.build}

                    **Why this matters:** This lab trains a practical part of the agentic AI loop. The important habit is to inspect evidence, not to assume that a fluent model output is correct.

                    **Concepts**

                    {md_table(lab.concepts)}

                    **Step-by-step**

                    {chr(10).join(f"{i}. {step}" for i, step in enumerate(lab.steps, 1))}

                    **Checkpoint**

                    {chr(10).join(f"- {item}" for item in lab.verify)}

                    **Trainer facilitation notes**

                    - Ask learners to show the exact execution or output that proves completion.
                    - Ask one learner to run an edge case while another observes the trace.
                    - Ask learners what they changed after evaluation and why.
                    - Do not accept a screenshot alone if the lab requires a workflow export, scorecard, transcript, or generated media.

                    **Common errors**

                    | Error | Likely cause | Fix |
                    |---|---|---|
                    {chr(10).join(f"| {a} | {b} | {c} |" for a, b, c in lab.troubleshooting)}

                    **Deliverable:** {lab.deliverable}
                    """
                ).strip()
            )
            shots = "\n\n".join(
                md for md in (shot_md(key, img_prefix) for key in lab_shots(lab)) if md
            )
            if shots:
                lab_sections.append("\n\n" + shots)
            # The turn-by-turn call belongs right after the voice booking lab.
            if lab.lab_no == "4":
                lab_sections.append("\n\n" + voice_conversation_section())

    appendices = dedent(
        """

        ## How you are assessed

        The deck tells you the appeal route is in this guide. Here it is, along with everything
        else you are entitled to know before you sit the papers.

        | | Written Assessment (SAQ) | Case Study (CS) |
        |---|---|---|
        | What it tests | Knowledge - the concepts behind the labs | Application - judgement on a workplace scenario |
        | Length | 6 open-ended questions | 6 scenario tasks |
        | Time | 60 minutes | 60 minutes |
        | Conditions | Open book, individually attempted | Open book, individually attempted |

        Both papers are sat on the final day, after an assessment briefing. There are no multiple
        choice questions. Nothing is assessed that the slides and the labs did not teach.

        **Marking.** Each criterion is graded **Competent (C)** or **Not Yet Competent (NYC)**. You
        are competent when every criterion on the paper is met.

        **If you are Not Yet Competent.** You are re-assessed **once**, on the failed instrument
        only. The trainer records the gap, you close it, and you sit that paper again. A
        re-assessment is a normal part of competency-based training - it is not a failure of the
        course.

        **Submission.** Complete your answers on the document provided and upload them to the LMS
        at https://lms-tms.tertiaryinfotech.com/.

        **Appeals.** If you believe an assessment decision is wrong, tell the trainer on the day and
        ask for the appeal form. State what you were marked NYC on and why you disagree. The
        assessment is reviewed by a second assessor, and you are told the outcome in writing. Raising
        an appeal never counts against you.

        ## Preparing for the Case Study

        The Case Study puts you back inside the four businesses you already built for - one task per business, each asking for the workflow, the design decisions and the trade-offs. Keep your workflow exports from the labs; the paper asks you to paste them.

        | Business | What they want | You built it in |
        |---|---|---|
        | Cook & Bake Academy | A website chat agent grounded ONLY in the course brochures, honest about what it does not know | Lab 3 |
        | MediRefill | A voice FAQ assistant that never gives medical advice - refusal plus pharmacist callback | Lab 5 |
        | GG Hair Salon | A voice receptionist that checks the REAL Google Calendar and books into it | Lab 4 |
        | GG News Studio | A presenter video from the day's facts, plus a cloud-vs-open-source production recommendation | Labs 7 and 7-os |

        For each answer, bring the same three things every lab demanded: **the workflow export**, **the execution trace that proves it ran**, and **one failure case you deliberately provoked** - the refusal, the taken slot, the question outside the PDF.

        ### What a strong answer includes

        - The n8n workflow JSON, exported after your own changes.
        - The nodes that matter, named and explained - not a node-by-node tour.
        - The guardrail and where it lives (prompt wording, credential boundary, refusal branch).
        - The trade-off, argued from what you saw in the labs, not from a vendor page.

        ### Assessment rubric

        | Criterion | Meets standard |
        |---|---|
        | Problem definition | Clear user, trigger, output, success criteria, and non-goals. |
        | Workflow correctness | Nodes are connected logically and executions prove the expected path. |
        | AI quality | Prompts, retrieval, and generated outputs are tested against a rubric. |
        | Guardrails | Secrets, personal data, unsupported claims, and publishing actions are controlled. |
        | Human review | High-risk decisions have approval or escalation. |
        | Documentation | Another operator can run, test, and recover the workflow. |

        ## Troubleshooting

        | Symptom | First check | Typical fix |
        |---|---|---|
        | n8n cannot reach Ollama | Ollama credential base URL | Use `http://host.docker.internal:11434`. |
        | Agent produces empty answer | Tool name, model output, execution trace | Rename tools simply and inspect model node output. |
        | RAG answer is invented | Prompt and retrieval trace | Add evidence-only instruction and refusal rule. |
        | Website cannot call n8n | Workflow active state and webhook URL | Activate workflow and use production webhook URL. |
        | Voice call fails | API key, microphone permission, session response | Fix credential and browser permission. |
        | Video generation fails | API quota, payload length, credential | Run a short test and inspect provider response. |
        | The generated video does not match the idea | The gemma4 shot script | Tighten the subject, camera and lighting wording, then regenerate. |

        ## Glossary

        - **Agent:** A workflow component that uses a model plus context, memory, or tools to complete a task.
        - **Agentic AI loop:** The engineering cycle used to design, test, improve, and operate AI workflows.
        - **Embedding:** A numeric representation of text used for similarity search.
        - **Grounding:** Requiring answers to come from approved sources or retrieved documents.
        - **Guardrail:** A rule, branch, approval, or technical boundary that reduces unsafe behavior.
        - **Human-in-the-loop:** A workflow design where a person reviews or approves high-risk actions.
        - **RAG:** Retrieval-Augmented Generation, where the model retrieves relevant information before answering.
        - **Tool call:** A controlled action the agent can request, such as creating a booking or searching a vector store.
        - **Vector store:** A database or memory store that retrieves text chunks by semantic similarity.
        - **Webhook:** A URL that starts a workflow when another app or browser sends an HTTP request.
        """
    ).strip()
    return normalize_md("\n".join(toc) + "\n" + intro + "".join(lab_sections) + "\n\n" + appendices)


DAYS = 2
ASSESSMENT_BLOCK = [
    ("Assessment briefing - instruments, timing, and what evidence is graded", 15),
    ("**Written Assessment (SAQ)** - 6 open-ended knowledge questions (K1-K6)", 60),
    ("**Case Study (CS)** - 6 scenario tasks drawn from the labs (A1-A6)", 60),
]
ASSESSMENT_MINUTES = sum(m for _, m in ASSESSMENT_BLOCK)


def _day_plan():
    """The fixed two-day shape of this course.

    Day 1 morning is Topic 1 (Chatbot), Day 1 afternoon is Topic 2 (Voice Agent),
    and Day 2 is Topic 3 (Video Agent) with the assessment at the end of the
    afternoon. 9:30-6:30 with a 1-hour lunch = 8 taught hours a day; tea breaks
    are counted inside the sessions, not on top of them.
    """
    day1 = [lab for lab in LABS if lab.topic in (1, 2)]
    day2 = [lab for lab in LABS if lab.topic == 3]
    day_minutes = 8 * 60
    for d, (labs, cap) in enumerate(
        ((day1, day_minutes), (day2, day_minutes - ASSESSMENT_MINUTES)), 1
    ):
        total = sum(l.minutes for l in labs)
        assert total <= cap, f"Day {d} has {total} min of labs but only {cap} min free"
    return [day1, day2]


def lesson_plan(slide_map: dict[str, int] | None = None) -> str:
    slide_map = slide_map or {}
    days = _day_plan()

    def hhmm(total: int) -> str:
        """Wall-clock time, `total` being minutes since midnight."""
        h, m = divmod(total, 60)
        ampm = "am" if h < 12 else "pm"
        h12 = h if h <= 12 else h - 12
        return f"{h12}:{m:02d}{ampm}"

    day_sections = []
    for d, labs in enumerate(days, 1):
        is_last = d == len(days)

        # Build the day as an ordered list of (title, minutes, deck ref, evidence), then
        # walk a real wall clock over it. The old version derived the clock from a taught-
        # minutes counter and bolted +60 on once it passed 1:00pm, which stretched whichever
        # activity happened to straddle noon into a two-hour block and printed the LUNCH row
        # out of sequence beside it.
        items = []
        for lab in labs:
            slide = slide_map.get(lab.lab_no)
            items.append((
                f"Lab {lab.lab_no} - {lab.title}", lab.minutes,
                f"Slide {slide}" if slide else "-", lab.deliverable,
            ))

        taught = sum(l.minutes for l in labs) + (ASSESSMENT_MINUTES if is_last else 0)
        # Every training day totals EXACTLY 8 instructional hours. What the labs do not fill
        # is concept delivery and evidence review - the deck's concept, workflow and quality
        # slides are taught time too, and calling that "review" understated it.
        slack = 8 * 60 - taught
        if slack > 0:
            items.append((
                "Concept delivery, review and evidence check - the trainer teaches the "
                "concept, workflow and quality slides for the day's labs, then learners "
                "show their executions, transcripts and scorecards",
                slack, "-", "Trainer sign-off on the day's evidence",
            ))
            taught += slack

        if is_last:
            for title, mins in ASSESSMENT_BLOCK:
                items.append((title, mins, "-", "Completed answer script submitted on the LMS"))

        # Lunch is a fixed hour, taken at the activity boundary nearest 1:00pm. It never
        # splits an activity, so the clock in this table is the clock in the room.
        bounds, run = [], 0
        for _, mins, _, _ in items:
            run += mins
            bounds.append(run)
        target = 210                                  # 210 min after 9:30am = 1:00pm
        lunch_after = min(range(len(bounds)), key=lambda i: abs(bounds[i] - target))

        rows, t = [], 9 * 60 + 30                     # the day starts at 9:30am
        for i, (title, mins, ref, evidence) in enumerate(items):
            rows.append(f"| {hhmm(t)} - {hhmm(t + mins)} | {mins} min | {title} | {ref} | {evidence} |")
            t += mins
            if i == lunch_after:
                rows.append(f"| {hhmm(t)} - {hhmm(t + 60)} | 60 min | **LUNCH** | - | - |")
                t += 60

        day_sections.append(
            f"### Day {d}\n\n"
            "| Time | Duration | Topic / Activity | Deck | Evidence produced |\n"
            "|---|---|---|---|---|\n"
            + "\n".join(rows)
            + f"\n\n**Day {d} instructional time: {taught // 60}h {taught % 60:02d}m** "
            f"(9:30am-{hhmm(t)}, less a 1-hour lunch; tea breaks are taken inside the sessions).\n"
        )

    topic_rows = []
    for idx, (topic, desc) in enumerate(TOPICS, 1):
        labs = ", ".join(f"Lab {lab.lab_no}" for lab in LABS if lab.topic == idx)
        lo, _lo_text, ka = TOPIC_LO[idx]
        topic_rows.append(f"| {idx} | {topic} | {desc} | {lo} ({ka}) | {labs} |")

    total_min = sum(l.minutes for l in LABS)

    return normalize_md(dedent(
        f"""
        # Lesson Plan - {COURSE_TITLE}

        ## Course profile

        | | |
        |---|---|
        | Provider | {ORG} |
        | Course code | {COURSE_CODE} |
        | WSQ skill | {WSQ_TSC} |
        | Version | {VERSION} ({VERSION_DATE}) |
        | Mode | Instructor-led adult training, hands-on labs |
        | Duration | {len(days)} training days, 8 instructional hours each |
        | Practical ratio | At least 70 percent hands-on lab time |
        | Trainer | Dr Alfred Ang |

        ## Trainer strategy

        Demonstrate each concept with one small working example, then move quickly into learner practice. For every lab, ask learners to show the **evidence** - the n8n execution trace, the generated media, the call transcript, the scorecard. Do not accept a screenshot of a finished screen as proof that the workflow is correct: a fluent output and a correct workflow are not the same thing, and telling them apart is the skill this course teaches.

        ## Topic map

        Each topic delivers one WSQ learning outcome of {WSQ_TSC}. The Written Assessment (SAQ) tests K1-K6; the Case Study tests A1-A6.

        | Session | Topic | Focus | WSQ outcome | Labs |
        |---|---|---|---|---|
        {chr(10).join(topic_rows)}

        ## Daily schedule

        Total taught content: {total_min // 60} hours {total_min % 60:02d} minutes across {len(LABS)} labs.

        {chr(10).join(day_sections)}

        ## Assessment

        | Instrument | Duration | What it tests | When |
        |---|---|---|---|
        | Briefing for Assessment | 15 min | The learner knows the instruments, the criteria and the appeal route | Before the assessment, on the final day |
        | Written Assessment (WA) - Short Answer | 60 min | KNOWLEDGE: the concepts behind the labs | Final day, after the briefing |
        | Case Study (CS) | 60 min | APPLICATION: judgement on a workplace scenario built from the labs | Final day, after the WA |

        Both instruments are **open book**, individually attempted, and marked **Competent / Not Yet Competent**. A learner assessed NYC is re-assessed once, on the failed instrument only.

        ## Evidence the trainer collects

        - Workflow exports and n8n execution traces for each lab.
        - The provoked failures: the RAG refusal (Lab 2), the course that does not exist (Lab 3), Ava's medical refusals (Lab 5).
        - The Google Calendar booking created by voice, with its call transcript (Lab 4).
        - The lip-sync comparison scorecard: MuseTalk vs HeyGen on one script (Lab 6).
        - Generated media: the news avatar videos (Labs 7 and 7-os), the Aria latency screenshot (Lab 8), the Nova session scorecard (Lab 9), the Veo clip and its shot prompt (Lab 10).
        """
    ))


def readme() -> str:
    lab_rows = "\n".join(
        f"| {lab.lab_no} | {lab.title} | {lab.minutes} min | [`labs/topic-{lab.topic:02d}/lab-{lab.lab_no}-{lab.slug}/`](labs/topic-{lab.topic:02d}/lab-{lab.lab_no}-{lab.slug}/) |"
        for lab in LABS
    )
    topic_rows = "\n".join(f"- **{title}:** {desc}" for title, desc in TOPICS)
    return normalize_md(dedent(
        f"""
        # Automate Video and Voice AI Agents with n8n

        **WSQ course - Tertiary Infotech Academy Pte Ltd**

        This repository contains high-quality adult-training courseware for building AI agents, RAG assistants, voice agents, and AI avatar/video automations with n8n.

        The rebuilt courseware follows an **agentic AI loop engineering** standard:

        **Define -> Build -> Observe -> Evaluate -> Improve -> Guardrail -> Document**

        ## Courseware deliverables

        - [Learner Guide Markdown](courseware/LG-Automate-Video-and-Voice-AI-Agents-with-n8n.md)
        - [Learner Guide DOCX](courseware/LG-Automate-Video-and-Voice-AI-Agents-with-n8n.docx)
        - [PowerPoint deck](courseware/{PPT_STEM}.pptx)
        - [Lesson Plan](courseware/LP-Automate-Video-and-Voice-AI-Agents-with-n8n.md)

        ## Topics

        {topic_rows}

        ## Labs

        | Lab | Title | Time | Folder |
        |---|---:|---:|---|
        {lab_rows}

        ## Existing runnable assets

        The original runnable n8n workflows and websites remain in `lab0/` to `lab10/` and `lab7-opensource/`. The new `labs/` folder provides the adult-training lab instructions and quality standard around those assets.

        ## Quick start

        ```bash
        cd lab0
        docker compose up -d
        ollama pull gemma4
        ollama pull nomic-embed-text
        ```

        Then open n8n at `http://localhost:5678` and create an Ollama credential with:

        ```text
        http://host.docker.internal:11434
        ```

        ## Security rule

        API keys belong in n8n credentials or local environment files. They must not be placed in browser JavaScript, Markdown examples, screenshots, workflow notes, or committed JSON.
        """
    ))


IMG_RE = re.compile(r"^!\[(?P<alt>[^\]]*)\]\((?P<src>[^)]+)\)\s*$")


def plain(text: str) -> str:
    """Strip Markdown emphasis/code marks. Slides render text literally, so a step
    written as **Publish** must not appear on the slide as asterisks."""
    text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
    text = re.sub(r"`([^`]+)`", r"\1", text)
    text = re.sub(r"(?<!\w)\*(.+?)\*(?!\w)", r"\1", text)
    return text.replace("**", "").replace("`", "")


def _prodoc():
    """Load prodoc.py (the WSQ house document helpers) by path."""
    import importlib.util

    spec = importlib.util.spec_from_file_location("prodoc", PRODOC)
    mod = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(mod)
    return mod


def _toc_field(doc, entries):
    """A real Word TOC field whose CACHED RESULT is the actual table of contents.

    Word recomputes the field when the document opens. LibreOffice - which is what turns
    these into the PDFs learners are given - does not: it renders whatever result the field
    already carries. With prodoc's placeholder cached in there, every shipped PDF printed
    "Right-click and choose Update Field..." where its contents page should be.

    So we cache a computed result. A field may span paragraphs: begin/instrText/separate,
    then the entry paragraphs, then end. Word still replaces all of it on update, so the
    DOCX stays a live TOC while the PDF finally shows one.
    """
    from docx.shared import Pt, Inches
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement
    from docx.enum.text import WD_TAB_ALIGNMENT, WD_TAB_LEADER

    head = doc.add_paragraph()
    r = head.add_run("TABLE OF CONTENTS")
    r.bold = True
    r.font.size = Pt(12)

    def fld(p, kind_):
        e = OxmlElement("w:fldChar")
        e.set(qn("w:fldCharType"), kind_)
        run = p.add_run()
        run._r.append(e)

    p0 = doc.add_paragraph()
    fld(p0, "begin")
    instr = OxmlElement("w:instrText")
    instr.set(qn("xml:space"), "preserve")
    instr.text = 'TOC \\o "1-3" \\h \\z \\u'
    p0.add_run()._r.append(instr)
    fld(p0, "separate")

    for level, text, page in entries:
        p = doc.add_paragraph()
        pf = p.paragraph_format
        pf.left_indent = Inches(0.25 * (level - 1))
        pf.space_after = Pt(2)
        pf.tab_stops.add_tab_stop(Inches(6.3), WD_TAB_ALIGNMENT.RIGHT, WD_TAB_LEADER.DOTS)
        run = p.add_run(f"{text}\t{page}")
        run.font.size = Pt(10 if level > 1 else 11)
        run.bold = level == 1

    pend = doc.add_paragraph()
    fld(pend, "end")


def write_docx(md: str, path: Path, kind: str = "Learner Guide", toc_entries=None) -> None:
    """Render a Markdown document as a WSQ house DOCX: cover page, Document Version
    Control Record, a real Word TOC field, Arial 11pt body and a Page X of Y footer."""
    try:
        from docx import Document
        from docx.shared import Inches, Pt
        from docx.enum.text import WD_BREAK
    except Exception as exc:  # pragma: no cover
        print(f"python-docx unavailable, skipped DOCX: {exc}")
        return

    pd = _prodoc()
    doc = Document()
    doc.styles["Normal"].font.name = "Arial"
    doc.styles["Normal"].font.size = Pt(11)
    pd.style_headings(doc)

    pd.add_cover_page(
        doc, kind, COURSE_TITLE, VERSION,
        org_logo=str(ORG_LOGO) if ORG_LOGO.exists() else None,
        course_logo=str(COURSE_LOGO) if COURSE_LOGO.exists() else None,
        course_code=COURSE_CODE,
    )
    # add_version_control already ends with a page break; a second one here left a blank
    # page sitting between the version record and the contents page.
    pd.add_version_control(doc, VERSION_ROWS)
    if toc_entries:
        _toc_field(doc, toc_entries)
        doc.add_paragraph().add_run().add_break(WD_BREAK.PAGE)
    else:
        pd.add_toc(doc)          # first pass: no page numbers yet; it breaks the page itself

    # The cover already carries the title block, so drop the markdown's own H1/meta.
    lines = md.splitlines()
    body = []
    seen_h1 = False
    for line in lines:
        if line.startswith("# ") and not seen_h1:
            seen_h1 = True
            continue
        if not body and (line.startswith("**Course code:**") or not line.strip()):
            continue
        body.append(line)

    in_table = False
    for line in body:
        img = IMG_RE.match(line)
        if img:
            src = (ROOT / img.group("src")).resolve()
            if src.exists():
                doc.add_picture(str(src), width=Inches(6.1))
                cap = doc.add_paragraph(img.group("alt"))
                cap.runs[0].italic = True
                cap.runs[0].font.size = Pt(9)
            continue

        if line.startswith("|"):
            # Markdown tables become real Word tables, not pipe soup.
            cells = [c.strip() for c in line.strip("|").split("|")]
            if set("".join(cells)) <= set("-: "):
                continue  # the |---|---| separator row
            if not in_table:
                tbl = doc.add_table(rows=0, cols=len(cells))
                tbl.style = "Light Grid Accent 1"
                in_table = True
            row = tbl.add_row().cells
            for i, c in enumerate(cells[: len(row)]):
                row[i].text = plain(c)
            continue
        in_table = False

        if line.startswith("#### "):
            doc.add_heading(plain(line[5:]), 3)
        elif line.startswith("### "):
            doc.add_heading(plain(line[4:]), 2)
        elif line.startswith("## "):
            doc.add_heading(plain(line[3:]), 1)
        elif line.startswith("- "):
            doc.add_paragraph(plain(line[2:]), style="List Bullet")
        elif line and line[0].isdigit() and ". " in line[:4]:
            # Word's "List Number" style shares ONE counter across the whole document, so
            # the steps ran on and on: Lab 1.2 opened at step 14, Lab 4.7 at 149. The
            # Markdown already numbers each lab from 1, so print ITS number as text and let
            # the two documents say the same thing.
            num, text = line.split(". ", 1)
            p = doc.add_paragraph(f"{num}. {plain(text)}")
            p.paragraph_format.left_indent = Inches(0.25)
        elif line.strip() == "":
            doc.add_paragraph("")
        else:
            doc.add_paragraph(plain(line))

    pd.add_page_numbers(doc, left_text=f"{COURSE_CODE} - {COURSE_TITLE}")
    pd.enable_update_fields(doc)
    doc.save(path)


SLIDE_MAP: dict[str, int] = {}   # lab_no -> the slide the lab starts on


def write_pptx(path: Path) -> dict[str, int]:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    BLUE = "1F6FEB"
    TEAL = "10B981"
    INK = "161B26"
    GREY = "5B6372"
    VIOLET = "7C3AED"
    PALE = "F5F8FC"
    LINE = "DCE5F0"
    AMBER = "F59E0B"
    RED = "E25555"
    WHITE = "FFFFFF"

    def rgb(value: str):
        return RGBColor.from_string(value)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def text_box(slide, text, x, y, w, h, size=18, color=INK, bold=False,
                 align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, margin=0.04):
        shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        shape.text_frame.clear()
        shape.text_frame.word_wrap = True
        shape.text_frame.margin_left = Inches(margin)
        shape.text_frame.margin_right = Inches(margin)
        shape.text_frame.margin_top = Inches(margin)
        shape.text_frame.margin_bottom = Inches(margin)
        shape.text_frame.vertical_anchor = valign
        p = shape.text_frame.paragraphs[0]
        p.text = str(text)
        p.alignment = align
        p.font.name = "Arial"
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = rgb(color)
        return shape

    def shape_box(slide, x, y, w, h, fill=PALE, radius=True, line=LINE, line_width=1):
        kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
        shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill)
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(line_width)
        return shape

    def circle(slide, label, x, y, d, fill=BLUE, color=WHITE, size=16):
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill)
        shape.line.fill.background()
        text_box(slide, label, x, y, d, d, size, color, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
        return shape

    def add_footer(slide, slide_no):
        shape_box(slide, 0.55, 7.03, 12.23, 0.015, LINE, False, LINE, 0)
        text_box(slide, "TGS-2024052081 - Automate Video and Voice AI Agents with n8n", 0.6, 7.09, 5.7, 0.22, 8, GREY)
        text_box(slide, "Copyright 2026 Tertiary Infotech Academy Pte Ltd", 4.65, 7.09, 4.1, 0.22, 8, GREY, False, PP_ALIGN.CENTER)
        text_box(slide, f"{slide_no:03d}", 12.15, 7.09, 0.55, 0.22, 8, GREY, True, PP_ALIGN.RIGHT)

    def base(title, kicker="COURSEWARE"):
        s = prs.slides.add_slide(blank)
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = rgb(WHITE)
        # The kicker gets the full slide width. At 3.0" a long one ("LAB 6.2 - BUILD AN
        # INTERACTIVE AVATAR THAT RENDERS IN THE BROWSER") wrapped onto three lines and
        # printed straight through the title underneath it.
        text_box(s, kicker.upper(), 0.62, 0.33, 12.0, 0.28, 11, BLUE, True)
        text_box(s, title, 0.62, 0.68, 12.0, 0.62, 28, INK, True)
        return s

    def card(slide, x, y, w, h, label, body, accent=BLUE, badge=None):
        shape_box(slide, x, y, w, h, WHITE, True, LINE, 1)
        shape_box(slide, x, y, 0.08, h, accent, False, accent, 0)
        if badge:
            circle(slide, badge, x + 0.26, y + 0.24, 0.48, accent, WHITE, 12)
            tx = x + 0.88
            tw = w - 1.12
        else:
            tx = x + 0.3
            tw = w - 0.55
        compact = h < 1.3
        text_box(slide, label, tx, y + (0.15 if compact else 0.22), tw, 0.34, 16 if compact else 17, INK, True)
        body_x = tx if compact and badge else x + 0.3
        body_w = tw if compact and badge else w - 0.55
        text_box(slide, body, body_x, y + (0.55 if compact else 0.72), body_w,
                 h - (0.62 if compact else 0.9), 11 if compact else 14, GREY)

    def topic_slide(index, title, description):
        s = base(title, f"TOPIC {index:02d}")
        text_box(s, f"{index:02d}", 9.9, 0.95, 2.5, 1.65, 74, "E7EFFA", True, PP_ALIGN.RIGHT)
        shape_box(s, 0.65, 1.65, 7.75, 2.0, PALE, True, PALE, 0)
        text_box(s, description, 1.0, 2.03, 6.95, 1.2, 20, INK, False, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
        # Seven steps have to fit LEFT of the LEARNING ARC card (which starts at x=9.15).
        # At the old 1.76" pitch, steps 6 and 7 - Guardrail and Document - ran underneath
        # it, so the loop this whole course is built on silently showed only five steps.
        labels = ["Define", "Build", "Observe", "Evaluate", "Improve", "Guardrail", "Document"]
        pitch = 1.17
        for i, label in enumerate(labels):
            x = 0.72 + i * pitch
            circle(s, str(i + 1), x, 4.38, 0.48, BLUE if i < 4 else TEAL, WHITE, 12)
            text_box(s, label, x - 0.31, 5.0, 1.10, 0.32, 11, INK, True, PP_ALIGN.CENTER)
            if i < len(labels) - 1:
                shape_box(s, x + 0.53, 4.59, pitch - 0.55, 0.045, LINE, False, LINE, 0)
        text_box(s, "Every lab produces visible evidence for trainer review.", 0.72, 5.72, 7.5, 0.5, 17, TEAL, True)
        shape_box(s, 9.15, 3.18, 3.1, 2.7, WHITE, True, LINE, 1)
        text_box(s, "LEARNING ARC", 9.5, 3.52, 2.4, 0.28, 11, VIOLET, True)
        text_box(s, "Concept", 9.5, 4.05, 2.3, 0.3, 18, INK, True)
        text_box(s, "Practice", 9.5, 4.52, 2.3, 0.3, 18, INK, True)
        text_box(s, "Evidence", 9.5, 4.99, 2.3, 0.3, 18, INK, True)
        text_box(s, "Reflection", 9.5, 5.46, 2.3, 0.3, 18, INK, True)

    def lab_target(lab):
        s = base(f"Lab {lab.lab_no}: {lab.title}", "BUILD TARGET")
        shape_box(s, 0.65, 1.55, 8.05, 3.55, PALE, True, PALE, 0)
        circle(s, lab.lab_no, 1.0, 1.92, 1.18, BLUE, WHITE, 21)
        text_box(s, "YOU WILL BUILD", 2.55, 1.88, 2.4, 0.3, 11, BLUE, True)
        # The box is middle-anchored, so text longer than it fits grows in BOTH directions
        # and prints over the two labels that bracket it. Size the type to the sentence:
        # roughly 30 characters fit on a 5.55" line at 24pt, and four lines fill the box.
        build_pt = 24 if len(lab.build) <= 95 else 20 if len(lab.build) <= 130 else 17
        text_box(s, lab.build, 2.55, 2.24, 5.55, 1.44, build_pt, INK, True, valign=MSO_ANCHOR.MIDDLE)
        text_box(s, "DELIVERABLE", 2.55, 3.8, 2.1, 0.28, 11, TEAL, True)
        text_box(s, lab.deliverable, 2.55, 4.18, 5.55, 0.62, 15, GREY)
        card(s, 9.08, 1.55, 3.6, 1.05, "TIMEBOX", f"{lab.minutes} minutes", AMBER, "T")
        card(s, 9.08, 2.85, 3.6, 1.05, "MODE", "Demo, guided build, independent test", VIOLET, "M")
        card(s, 9.08, 4.15, 3.6, 1.05, "EVIDENCE", "Execution trace plus learner artifact", TEAL, "E")
        text_box(s, "Success is demonstrated, not assumed.", 0.72, 5.72, 8.0, 0.5, 20, VIOLET, True)

    def concept_slide(lab):
        s = base(f"Lab {lab.lab_no}: Four Ideas That Make It Work", "CONCEPT MAP")
        colors = [BLUE, TEAL, VIOLET, AMBER]
        positions = [(0.68, 1.52), (6.78, 1.52), (0.68, 4.02), (6.78, 4.02)]
        for i, ((name, desc), (x, y)) in enumerate(zip(lab.concepts, positions)):
            card(s, x, y, 5.85, 2.08, plain(name), plain(desc), colors[i], str(i + 1))
        text_box(s, "Connect the four ideas before touching the workflow canvas.", 0.75, 6.35, 8.8, 0.38, 16, GREY)

    def step_slide(lab):
        """Build-sequence slides. Steps are PAGINATED (six per slide) so a long lab
        never silently loses steps, and the font shrinks to fit the longest step."""
        per_slide = 6
        pages = [lab.steps[i:i + per_slide] for i in range(0, len(lab.steps), per_slide)] or [[]]

        for page_no, page in enumerate(pages, 1):
            suffix = "" if len(pages) == 1 else f" ({page_no}/{len(pages)})"
            s = base(f"Lab {lab.lab_no}: Build Sequence{suffix}", "WORKFLOW PLAN")

            # One font size for the whole page, driven by its longest step, so the
            # boxes stay a uniform grid instead of each one sizing itself.
            longest = max((len(plain(step)) for step in page), default=0)
            if longest <= 105:
                size = 14
            elif longest <= 150:
                size = 12
            elif longest <= 210:
                size = 11
            elif longest <= 280:
                size = 10
            else:
                size = 9

            for i, step in enumerate(page):
                col = i % 2
                row = i // 2
                x = 0.72 + col * 6.25
                y = 1.45 + row * 1.78
                n = (page_no - 1) * per_slide + i + 1
                circle(s, str(n), x, y + 0.15, 0.58, BLUE if col == 0 else TEAL, WHITE, 14)
                shape_box(s, x + 0.82, y, 5.13, 1.5, WHITE, True, LINE, 1)
                text_box(s, plain(step), x + 1.05, y + 0.1, 4.68, 1.3, size, INK, False,
                         valign=MSO_ANCHOR.MIDDLE)

            text_box(s, "Build one observable behavior at a time.", 0.75, 6.62, 6.5, 0.3, 15, VIOLET, True)

    def verify_slide(lab):
        s = base(f"Lab {lab.lab_no}: Prove, Evaluate, Improve", "QUALITY GATE")
        shape_box(s, 0.68, 1.48, 7.65, 4.92, PALE, True, PALE, 0)
        checks = lab.verify + ["Run one edge case and record the change made after evaluation."]
        checks = checks[:6]
        # Spacing and font follow the row count so a 6-check lab still fits the panel.
        gap = 0.91 if len(checks) <= 5 else 0.78
        size = 15 if max((len(plain(c)) for c in checks), default=0) <= 95 else 13
        for i, item in enumerate(checks):
            y = 1.82 + i * gap
            circle(s, "OK", 1.02, y, 0.5, TEAL, WHITE, 10)
            text_box(s, plain(item), 1.78, y - 0.02, 5.95, gap - 0.06, size, INK, valign=MSO_ANCHOR.MIDDLE)
        shape_box(s, 8.72, 1.48, 3.95, 4.92, WHITE, True, LINE, 1)
        text_box(s, "EVIDENCE STACK", 9.08, 1.84, 3.2, 0.3, 11, BLUE, True)
        evidence = [("01", "Input", "Exact test payload"), ("02", "Trace", "Node execution path"), ("03", "Output", "Generated result"), ("04", "Decision", "Rubric and revision")]
        for i, (num, label, note) in enumerate(evidence):
            y = 2.38 + i * 0.87
            circle(s, num, 9.08, y, 0.45, VIOLET, WHITE, 10)
            text_box(s, label, 9.72, y - 0.02, 1.0, 0.25, 14, INK, True)
            text_box(s, note, 10.72, y - 0.02, 1.48, 0.45, 12, GREY)
        text_box(s, "PASS", 9.08, 5.92, 1.0, 0.28, 12, TEAL, True)
        shape_box(s, 10.12, 5.97, 2.12, 0.12, TEAL, False, TEAL, 0)

    def fixes_slide(lab):
        """Troubleshooting table, PAGINATED three rows per slide so a lab with five
        known failure modes does not quietly show only the first three."""
        rows = lab.troubleshooting
        pages = [rows[i:i + 3] for i in range(0, len(rows), 3)] or [[]]
        headers = [("SYMPTOM", RED), ("LIKELY CAUSE", AMBER), ("TARGETED FIX", TEAL)]
        xs = [0.68, 4.72, 8.76]

        for page_no, page in enumerate(pages, 1):
            suffix = "" if len(pages) == 1 else f" ({page_no}/{len(pages)})"
            s = base(f"Lab {lab.lab_no}: Diagnose Before Rebuilding{suffix}", "TROUBLESHOOTING")
            for (label, color), x in zip(headers, xs):
                shape_box(s, x, 1.47, 3.78, 0.48, color, True, color, 0)
                text_box(s, label, x, 1.54, 3.78, 0.25, 11, WHITE, True, PP_ALIGN.CENTER)

            longest = max((len(plain(v)) for row in page for v in row), default=0)
            size = 14 if longest <= 60 else (12 if longest <= 95 else 11)

            for row_i, (err, cause, fix) in enumerate(page):
                y = 2.18 + row_i * 1.43
                colors = ["FFF4F4", "FFF8E8", "ECFDF5"]
                for col, value in enumerate([err, cause, fix]):
                    shape_box(s, xs[col], y, 3.78, 1.12, colors[col], True, colors[col], 0)
                    text_box(s, plain(value), xs[col] + 0.18, y + 0.1, 3.42, 0.92, size, INK,
                             col == 0, valign=MSO_ANCHOR.MIDDLE)
            text_box(s, "Change one variable, rerun the same test, compare the trace.", 0.75, 6.55, 9.5, 0.3, 16, VIOLET, True)

    s = prs.slides.add_slide(blank)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = rgb(WHITE)
    # The cover carries the house identity as LOGOS, the way the LP and LG covers do -
    # a text-only cover is not the WSQ cover page.
    if ORG_LOGO.exists():
        s.shapes.add_picture(str(ORG_LOGO), Inches(0.72), Inches(0.42), height=Inches(0.52))
    if COURSE_LOGO.exists():
        # Beside the org logo, NOT top-right: the accent panel is added after this and
        # would paint straight over it (in pptx, later shapes win).
        s.shapes.add_picture(str(COURSE_LOGO), Inches(1.62), Inches(0.44), height=Inches(0.48))
    text_box(s, "TERTIARY INFOTECH ACADEMY", 0.72, 1.02, 4.5, 0.32, 12, BLUE, True)
    text_box(s, "Automate Video and Voice\nAI Agents with n8n", 0.72, 1.62, 7.35, 1.9, 38, INK, True, valign=MSO_ANCHOR.MIDDLE)
    text_box(s, "Agentic AI loop engineering for adult workplace training", 0.78, 3.72, 6.75, 0.72, 20, GREY)
    text_box(s, f"{COURSE_CODE}  |  {VERSION}  ·  {VERSION_DATE}  |  {UEN}",
             0.78, 5.55, 7.4, 0.38, 12, GREY, True)
    shape_box(s, 8.33, 0.62, 4.32, 5.95, PALE, True, PALE, 0)
    nodes = [(9.0, 1.2, "VOICE", BLUE), (10.7, 2.12, "n8n", VIOLET), (8.88, 3.25, "AI", TEAL), (10.78, 4.3, "VIDEO", AMBER), (9.02, 5.36, "REVIEW", RED)]
    for i, (x, y, label, color) in enumerate(nodes):
        if i < len(nodes) - 1:
            nx, ny, _, _ = nodes[i + 1]
            line = s.shapes.add_connector(1, Inches(x + 0.62), Inches(y + 0.62), Inches(nx + 0.62), Inches(ny + 0.62))
            line.line.color.rgb = rgb(LINE)
            line.line.width = Pt(3)
        circle(s, label, x, y, 1.22, color, WHITE, 13)


    # ── WSQ admin block ──────────────────────────────────────────────────────
    def traqom_slide(where: str) -> None:
        s = base("Digital Attendance - TRAQOM", "ADMIN")
        text_box(s, f"Please take attendance {where} of the session.", 0.75, 1.45, 11.9, 0.4, 19, GREY)
        for i, (n, label, note, color) in enumerate([
            ("1", "Open TRAQOM", "Scan the QR code shown by the trainer", BLUE),
            ("2", "Verify identity", "SingPass / NRIC as prompted", TEAL),
            ("3", "Confirm", "Attendance is recorded against TGS-2024052081", VIOLET),
        ]):
            y = 2.3 + i * 1.25
            circle(s, n, 1.0, y, 0.62, color, WHITE, 15)
            text_box(s, label, 1.95, y + 0.02, 3.2, 0.35, 18, INK, True)
            text_box(s, note, 5.3, y + 0.02, 7.0, 0.5, 15, GREY)
        shape_box(s, 0.72, 6.15, 11.9, 0.85, PALE, True, PALE, 0)
        text_box(s, "Attendance is a funding requirement. No attendance record, no course fee funding.",
                 1.05, 6.3, 11.3, 0.55, 15, RED, True, valign=MSO_ANCHOR.MIDDLE)

    traqom_slide("at the START")

    # Two trainer profiles, on two SEPARATE pages — the house checklist asks for a general
    # trainer template card AND the named trainer, each as its own visual profile page.
    def trainer_slide(kicker, title, badge, badge_color, name, creds, lines, footer, footer_color):
        s = base(title, kicker)
        shape_box(s, 3.55, 1.5, 6.3, 5.05, PALE if badge == "T" else WHITE, True, LINE, 1)
        circle(s, badge, 6.05, 1.95, 1.7, badge_color, WHITE, 40 if len(badge) == 1 else 30)
        text_box(s, name, 3.75, 3.9, 5.9, 0.45, 24, INK, True, PP_ALIGN.CENTER)
        text_box(s, creds, 3.75, 4.42, 5.9, 0.35, 15, GREY, False, PP_ALIGN.CENTER)
        text_box(s, lines, 3.95, 4.95, 5.5, 1.45, 15, GREY, False, PP_ALIGN.CENTER)
        text_box(s, footer, 3.75, 6.15, 5.9, 0.3, 12, footer_color, True, PP_ALIGN.CENTER)

    trainer_slide(
        "YOUR TRAINER", "About the Trainer", "?", BLUE,
        "<< Trainer Name >>", "<< Qualification >>  ·  << Years of experience >>",
        "<< Industry practice >>\n<< Teaching experience >>\n<< Specialisation >>",
        "General trainer template", BLUE,
    )
    trainer_slide(
        "YOUR TRAINER", "About the Trainer", "AA", VIOLET,
        "Dr Alfred Ang", "PhD  ·  ACTA / ACLP certified adult educator",
        "Founder, Tertiary Infotech Academy\nAI automation, voice and video agents\n"
        "n8n, RAG, LLM systems in production",
        "Your trainer today", VIOLET,
    )

    # Learner introductions — right after the trainer has introduced himself.
    s = base("Let's Get to Know Each Other", "NOW YOU - INTRODUCE YOURSELF")
    text_box(s, "One minute each. The trainer just introduced himself - your turn.", 0.75, 1.42, 11.9, 0.35, 17, GREY)
    intro_items = [
        ("1", "Your name and role", "And the organisation or team you work with", BLUE),
        ("2", "Your automation experience", "n8n, chatbots, voice or video AI - or none at all, which is fine", TEAL),
        ("3", "Why you are here", "The one thing you want to automate after these two days", VIOLET),
        ("4", "One fun fact", "Optional - break the ice", AMBER),
    ]
    for i, (n, label, note, color) in enumerate(intro_items):
        y = 2.0 + i * 1.02
        circle(s, n, 0.9, y, 0.6, color, WHITE, 15)
        text_box(s, label, 1.8, y + 0.02, 3.6, 0.4, 17, INK, True)
        text_box(s, note, 5.55, y + 0.04, 7.0, 0.5, 14, GREY)
    shape_box(s, 0.72, 6.25, 11.9, 0.75, PALE, True, PALE, 0)
    text_box(s, "You will be demoing your builds to each other for two days - start talking now.",
             1.05, 6.37, 11.3, 0.5, 15, VIOLET, True, valign=MSO_ANCHOR.MIDDLE)

    # Ground rules
    s = base("Ground Rules", "HOUSEKEEPING")
    rules = [
        ("Be on time", "Sessions start at 9:30am sharp", BLUE),
        ("Phones on silent", "Take calls outside the room", TEAL),
        ("Ask as you go", "A question now saves an hour later", VIOLET),
        ("Work the labs", "70% of this course is hands-on", AMBER),
        ("Respect the room", "One conversation at a time", RED),
    ]
    for i, (label, note, color) in enumerate(rules):
        y = 1.6 + i * 1.05
        circle(s, str(i + 1), 0.85, y, 0.58, color, WHITE, 14)
        text_box(s, label, 1.75, y - 0.02, 3.2, 0.35, 18, INK, True)
        text_box(s, note, 5.1, y - 0.02, 7.4, 0.45, 15, GREY)

    # Learning outcomes — the WSQ skill's three LOs, verbatim, with their K/A refs.
    s = base("Learning Outcomes", "WSQ SKILL - ARTIFICIAL INTELLIGENCE APPLICATION (AER-TEM-4026-1.1)")
    topic_for_lo = ["Topic 1 - Chatbot", "Topic 2 - Voice Agent", "Topic 3 - Video Agent"]
    for i, (lo, lo_text, ka) in enumerate(LEARNING_OUTCOMES):
        y = 1.62 + i * 1.42
        color = [BLUE, TEAL, VIOLET][i]
        circle(s, lo, 0.85, y + 0.18, 0.72, color, WHITE, 14)
        shape_box(s, 1.85, y, 10.75, 1.18, PALE, True, PALE, 0)
        text_box(s, lo_text, 2.15, y + 0.12, 10.15, 0.62, 16, INK, True)
        text_box(s, f"{ka}   ·   delivered through {topic_for_lo[i]}", 2.15, y + 0.78, 10.15, 0.3, 12, GREY)
    shape_box(s, 0.72, 6.15, 11.9, 0.82, PALE, True, PALE, 0)
    text_box(s, "The chatbots, voice agents and video avatars you build ARE the AI digital human applications this skill describes. "
                "The WA tests K1-K6; the Case Study tests A1-A6.",
             1.05, 6.28, 11.3, 0.6, 13, VIOLET, True, valign=MSO_ANCHOR.MIDDLE)

    # Lesson plan slide — the two-day shape the user actually experiences.
    s = base("Lesson Plan - 2 Days", "HOW THE DAYS RUN")
    text_box(s, "9:30am - 6:30pm  ·  1-hour lunch  ·  tea breaks inside the sessions  ·  8 instructional hours a day",
             0.75, 1.42, 11.9, 0.35, 16, GREY)
    sessions = [
        ("Day 1 · Morning", "Topic 1 - Chatbot", "Labs 0-3: local stack, first agent, RAG chatbot, CX advisor", BLUE),
        ("Day 1 · Afternoon", "Topic 2 - Voice Agent", "Labs 4-5: ElevenLabs booking agent, Vapi FAQ agent", TEAL),
        ("Day 2 · Morning", "Topic 3 - Video Agent", "Labs 6-10: lip-sync, news avatars, interactive avatars, Veo 3", VIOLET),
        ("Day 2 · Afternoon", "Topic 3 continues + review", "Evidence check, then WA (SAQ) 1h + Case Study 1h", AMBER),
    ]
    for i, (when, what, detail, color) in enumerate(sessions):
        y = 2.0 + i * 1.12
        shape_box(s, 0.72, y, 11.9, 0.95, PALE, True, PALE, 0)
        shape_box(s, 0.72, y, 0.09, 0.95, color, False, color, 0)
        text_box(s, when, 1.05, y + 0.1, 2.7, 0.35, 15, color, True)
        text_box(s, what, 1.05, y + 0.5, 2.9, 0.35, 13, INK, True)
        text_box(s, detail, 4.15, y + 0.1, 8.3, 0.75, 14, GREY, False, valign=MSO_ANCHOR.MIDDLE)

    # Download course material — a visual, not a bare link
    s = base("Download Your Course Material", "BEFORE WE START")
    text_box(s, "Everything - slides, Learner Guide, lab files - is on the LMS.", 0.75, 1.42, 11.9, 0.35, 17, GREY)
    steps = [
        ("1", "Go to", "lms-tms.tertiaryinfotech.com", BLUE),
        ("2", "Sign in", "Use the email you registered with", TEAL),
        ("3", "Open the course", "Automate Video and Voice AI Agents with n8n", VIOLET),
        ("4", "Download", "Learner Slides (PDF), Learner Guide (PDF), lab files", AMBER),
    ]
    for i, (n, label, note, color) in enumerate(steps):
        y = 2.0 + i * 1.15
        circle(s, n, 0.9, y, 0.6, color, WHITE, 15)
        text_box(s, label, 1.8, y, 2.6, 0.35, 17, INK, True)
        text_box(s, note, 4.5, y, 8.0, 0.5, 15, GREY)
    shape_box(s, 0.72, 6.3, 11.9, 0.72, PALE, True, PALE, 0)
    text_box(s, "lms-tms.tertiaryinfotech.com", 1.05, 6.4, 11.3, 0.5, 18, BLUE, True, valign=MSO_ANCHOR.MIDDLE)

    # Download the labs from GitHub — the runnable flows and websites live in one repo.
    s = base("Download the Labs from GitHub", "GET THE LAB FILES")
    text_box(s, "Every lab - the n8n flows, the websites, the PDFs, the launchers - is one public repository.",
             0.75, 1.42, 11.9, 0.35, 17, GREY)
    shape_box(s, 0.72, 1.95, 11.9, 1.0, PALE, True, LINE, 1)
    text_box(s, GITHUB_URL.replace("https://", ""), 1.0, 2.22, 11.4, 0.5, 17, BLUE, True, PP_ALIGN.CENTER)
    gh_steps = [
        ("1", "Open the repository", "Scan the link above, or search 'tertiarycourses' on GitHub", BLUE),
        ("2", "Code -> Download ZIP", "Or, if you use git:  git clone " + GITHUB_URL.replace("https://", ""), TEAL),
        ("3", "Windows: Unblock the ZIP", "Right-click the ZIP -> Properties -> tick Unblock -> OK, THEN extract", RED),
        ("4", "Open labs_local_n8n/", "The local build used in class - each lab folder has the flow JSON + website", VIOLET),
    ]
    for i, (n, label, note, color) in enumerate(gh_steps):
        y = 3.25 + i * 0.82
        circle(s, n, 0.9, y, 0.56, color, WHITE, 14)
        text_box(s, label, 1.8, y + 0.02, 3.5, 0.4, 16, INK, True)
        text_box(s, note, 5.45, y + 0.04, 7.1, 0.5, 13, GREY)
    shape_box(s, 0.72, 6.42, 11.9, 0.52, PALE, True, PALE, 0)
    text_box(s, "Cannot run Docker?  labs_remote_n8n/ is the same ten labs on the hosted n8n - ask the trainer.",
             1.05, 6.48, 11.3, 0.4, 13, VIOLET, True, valign=MSO_ANCHOR.MIDDLE)

    # The assessment is explained UP FRONT as well as before the papers: briefing,
    # the two instruments, then the flow diagram - the house admin-block order.
    def briefing_slide():
        s = base("Briefing for Assessment", "BEFORE YOU ARE ASSESSED")
        text_box(s, "Read this before the papers are handed out. You are entitled to know exactly how you will be judged.",
                 0.75, 1.42, 11.9, 0.4, 17, GREY)
        briefing = [
            ("What you sit", "A Written Assessment (short answer) and a Case Study.", BLUE),
            ("How long", "60 minutes each. Open book. Individually attempted.", TEAL),
            ("How it is marked", "Competent / Not Yet Competent against the assessment criteria.", VIOLET),
            ("If you are NYC", "You are re-assessed once, on the failed instrument only.", AMBER),
            ("Appeals", "Tell the trainer on the day; the appeal route is in your Learner Guide.", RED),
        ]
        for i, (label, note, color) in enumerate(briefing):
            y = 2.05 + i * 0.95
            circle(s, str(i + 1), 0.85, y, 0.58, color, WHITE, 14)
            text_box(s, label, 1.75, y - 0.02, 3.1, 0.4, 17, INK, True, valign=MSO_ANCHOR.MIDDLE)
            text_box(s, note, 5.0, y - 0.02, 7.6, 0.5, 15, GREY, False, valign=MSO_ANCHOR.MIDDLE)

    def assessment_slide():
        s = base("Assessment", "HOW YOU ARE ASSESSED")
        for i, (name, dur, tests, color) in enumerate([
            ("Written Assessment (WA)", "60 min", "KNOWLEDGE - the concepts behind the labs: RAG, grounding, webhooks, lip-sync engines, guardrails. (K1-K6)", BLUE),
            ("Case Study (CS)", "60 min", "APPLICATION - judgement on the four businesses you built for in the labs. (A1-A6)", VIOLET),
        ]):
            x = 0.75 + i * 6.1
            shape_box(s, x, 1.6, 5.85, 4.2, PALE, True, LINE, 1)
            shape_box(s, x, 1.6, 5.85, 0.6, color, True, color, 0)
            text_box(s, name, x, 1.73, 5.85, 0.35, 16, WHITE, True, PP_ALIGN.CENTER)
            circle(s, dur, x + 2.35, 2.4, 1.15, color, WHITE, 15)
            text_box(s, tests, x + 0.35, 3.9, 5.15, 1.6, 15, INK, False, PP_ALIGN.CENTER)
        shape_box(s, 0.75, 6.05, 11.85, 0.95, PALE, True, PALE, 0)
        text_box(s, "Open book · individually attempted · marked Competent / Not Yet Competent · one re-assessment on the failed instrument",
                 1.05, 6.2, 11.3, 0.6, 15, VIOLET, True, PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    briefing_slide()
    assessment_slide()

    s = base("Assessment Flow", "FROM PAPER TO SIGN-OFF")
    for i, (label, note, color) in enumerate([
        ("BRIEFING", "criteria explained", BLUE),
        ("ATTENDANCE", "TRAQOM, digital", TEAL),
        ("WA (SAQ)", "60 min · knowledge", BLUE),
        ("CASE STUDY", "60 min · application", VIOLET),
        ("SUBMIT", "upload on the LMS", TEAL),
        ("MARKING", "C / NYC per criterion", AMBER),
        ("SIGN-OFF", "Assessment Summary Record", RED),
    ]):
        x = 0.62 + i * 1.82
        shape_box(s, x, 2.3, 1.62, 1.75, WHITE, True, color, 1)
        shape_box(s, x, 2.3, 1.62, 0.44, color, True, color, 0)
        text_box(s, label, x, 2.37, 1.62, 0.3, 10, WHITE, True, PP_ALIGN.CENTER)
        text_box(s, note, x + 0.08, 2.88, 1.46, 1.0, 10, INK, False, PP_ALIGN.CENTER)
        if i < 6:
            shape_box(s, x + 1.66, 3.14, 0.42, 0.06, LINE, False, LINE, 0)
            text_box(s, ">", x + 1.76, 2.95, 0.3, 0.3, 15, LINE, True, PP_ALIGN.CENTER)
    shape_box(s, 0.62, 4.5, 12.0, 1.3, PALE, True, PALE, 0)
    text_box(s, "BOTH PAPERS ARE SAT ON DAY 2 AFTERNOON", 0.95, 4.65, 6.0, 0.3, 12, VIOLET, True)
    text_box(s, "Everything they ask comes from the labs you will have already built. Keep your workflow exports, traces and transcripts as you go.",
             0.95, 5.0, 11.4, 0.6, 15, INK, True)

    s = base("What You Will Be Able to Build", "COURSE OUTCOMES")
    outcomes = [
        ("LOCAL AI", "n8n agents powered by Ollama", BLUE),
        ("RAG", "Evidence-grounded workplace answers", TEAL),
        ("VOICE", "Secure browser-to-agent handoffs", VIOLET),
        ("VIDEO", "Avatar and cinematic automation", AMBER),
        ("OPERATIONS", "Tests, guardrails, review and recovery", RED),
    ]
    for i, (label, body, color) in enumerate(outcomes):
        x = 0.7 + (i % 3) * 4.15
        y = 1.55 + (i // 3) * 2.32
        w = 3.86 if i < 3 else 5.95
        if i >= 3:
            x = 0.7 + (i - 3) * 6.18
        card(s, x, y, w, 1.82, label, body, color, str(i + 1))
    text_box(s, "The course ends with an integrated, reviewable AI workforce capstone.", 0.75, 6.35, 10.8, 0.35, 16, GREY)

    # ---- Setup: the same stack on both operating systems ----
    s = base("Set Up Your Machine - macOS or Windows", "ENVIRONMENT")
    text_box(s, "Every lab runs on both. The tools are identical; only the install command and a little typing differ.", 0.75, 1.4, 11.9, 0.35, 17, GREY)

    shape_box(s, 0.72, 1.95, 5.9, 0.55, INK, True, INK, 0)
    text_box(s, "macOS  -  Terminal  +  Homebrew (brew.sh)", 0.72, 2.07, 5.9, 0.32, 14, WHITE, True, PP_ALIGN.CENTER)
    shape_box(s, 6.95, 1.95, 5.68, 0.55, BLUE, True, BLUE, 0)
    text_box(s, "Windows  -  PowerShell  +  winget (built in)", 6.95, 2.07, 5.68, 0.32, 14, WHITE, True, PP_ALIGN.CENTER)

    tools = [
        ("Docker Desktop", "brew install --cask docker", "winget install Docker.DockerDesktop"),
        ("Ollama", "brew install ollama", "winget install Ollama.Ollama"),
        ("Python 3", "already installed (python3)", "winget install Python.Python.3.12"),
        ("ffmpeg", "brew install ffmpeg", "winget install Gyan.FFmpeg"),
        ("ngrok  (Labs 4-5)", "brew install ngrok", "winget install ngrok.ngrok"),
    ]
    for i, (tool, mac, win) in enumerate(tools):
        y = 2.68 + i * 0.62
        fill = PALE if i % 2 == 0 else WHITE
        shape_box(s, 0.72, y, 11.9, 0.55, fill, False, LINE, 1)
        text_box(s, tool, 0.92, y + 0.05, 2.6, 0.45, 13, INK, True, valign=MSO_ANCHOR.MIDDLE)
        text_box(s, mac, 3.6, y + 0.05, 4.3, 0.45, 12, GREY, False, valign=MSO_ANCHOR.MIDDLE)
        text_box(s, win, 8.0, y + 0.05, 4.5, 0.45, 12, BLUE, False, valign=MSO_ANCHOR.MIDDLE)

    shape_box(s, 0.72, 5.9, 11.9, 1.0, PALE, True, PALE, 0)
    text_box(s, "WINDOWS GOTCHAS", 1.0, 6.0, 3.0, 0.28, 12, RED, True)
    text_box(s, "Reopen PowerShell after installing, or the command is 'not recognized'  ·  Docker needs WSL 2 (wsl --install)  ·  The command is python, NOT python3", 1.0, 6.3, 11.4, 0.5, 14, INK, True)

    s = base("The Agentic AI Loop", "ENGINEERING METHOD")
    loop = [
        ("1", "DEFINE", "Outcome + risk"), ("2", "BUILD", "Smallest proof"),
        ("3", "OBSERVE", "Trace behavior"), ("4", "EVALUATE", "Test rubric"),
        ("5", "IMPROVE", "One variable"), ("6", "GUARDRAIL", "Control action"),
        ("7", "DOCUMENT", "Make repeatable"),
    ]
    for i, (num, label, note) in enumerate(loop):
        x = 0.58 + i * 1.8
        color = BLUE if i < 3 else (VIOLET if i < 5 else TEAL)
        circle(s, num, x + 0.38, 2.02, 0.85, color, WHITE, 18)
        text_box(s, label, x, 3.13, 1.62, 0.3, 13, INK, True, PP_ALIGN.CENTER)
        text_box(s, note, x, 3.55, 1.62, 0.62, 12, GREY, False, PP_ALIGN.CENTER)
        if i < 6:
            shape_box(s, x + 1.32, 2.41, 0.5, 0.07, LINE, False, LINE, 0)
    shape_box(s, 1.02, 4.75, 11.28, 1.13, PALE, True, PALE, 0)
    text_box(s, "Learner habit", 1.35, 5.08, 1.55, 0.3, 13, VIOLET, True)
    text_box(s, "Show the input, trace, output, evaluation decision, and the next revision.", 3.05, 4.98, 8.75, 0.48, 20, INK, True, valign=MSO_ANCHOR.MIDDLE)

    def screenshot_slide(key: str, kicker: str) -> None:
        """One full-bleed screenshot slide, shown INSIDE the lab it belongs to."""
        if key not in SCREENSHOTS:
            return
        fname, caption = SCREENSHOTS[key]
        img = SHOTS_DIR / fname
        if not img.exists():
            print(f"  (skipped slide, screenshot not captured: {fname})")
            return

        # The title is the caption's opening clause. Split on ": " and NOT on a bare colon -
        # "http://localhost:8137" chopped one title to "Lab 5.1 - Digital Human Studio (http".
        # And a caption with no colon at all became a whole sentence, which wrapped down into
        # the screenshot below it. The image is pinned, so the title must stay one line.
        title = caption.split(": ", 1)[0].strip()
        if len(title) > 62:
            title = title.split(". ", 1)[0].strip().rstrip(".")
        if len(title) > 62:
            # Cut at a WORD boundary. A blind character slice left slide 98 reading
            # "...the Book by Voice call to acti".
            title = title[:62].rsplit(" ", 1)[0].rstrip(" ,-")
        s = base(title, kicker)

        # Fit the image inside the content area without distorting it.
        max_w, max_h = 11.6, 4.55
        try:
            from PIL import Image  # type: ignore

            with Image.open(img) as im:
                iw, ih = im.size
            ratio = min(max_w / iw, max_h / ih) if iw and ih else 0
            w, h = (iw * ratio, ih * ratio) if ratio else (max_w, max_h)
        except Exception:
            w, h = max_w, max_h * 0.9

        left = (13.333 - w) / 2
        top = 1.45 + (max_h - h) / 2
        shape_box(s, left - 0.06, top - 0.06, w + 0.12, h + 0.12, PALE, True, LINE, 1)
        s.shapes.add_picture(str(img), Inches(left), Inches(top), Inches(w), Inches(h))
        text_box(s, caption, 0.78, 6.35, 11.8, 0.5, 15, GREY, False, PP_ALIGN.CENTER)

    def arrow(slide, x, y, color=LINE):
        """Small horizontal connector between pipeline boxes."""
        shape_box(slide, x, y, 0.42, 0.06, color, False, color, 0)
        text_box(slide, ">", x + 0.1, y - 0.19, 0.3, 0.3, 15, color, True, PP_ALIGN.CENTER)

    def rag_concept_slides() -> None:
        # ---- 1. Why RAG: the hallucination problem, as a before/after ----
        s = base("Why RAG? The Model Has Never Seen Your Documents", "RAG - THE PROBLEM")
        shape_box(s, 0.7, 1.5, 5.85, 4.05, PALE, True, LINE, 1)
        shape_box(s, 0.7, 1.5, 5.85, 0.52, RED, True, RED, 0)
        text_box(s, "WITHOUT RAG", 0.7, 1.62, 5.85, 0.3, 13, WHITE, True, PP_ALIGN.CENTER)
        text_box(s, "“What is the refund policy?”", 1.05, 2.25, 5.15, 0.4, 17, INK, True)
        text_box(s, "The model was never trained on your handbook.\nIt does not say “I don't know”.", 1.05, 2.8, 5.15, 0.8, 15, GREY)
        shape_box(s, 1.05, 3.75, 5.15, 1.45, WHITE, True, RED, 1)
        text_box(s, "“Refunds are accepted within 30 days…”", 1.25, 3.95, 4.75, 0.5, 15, INK, True)
        text_box(s, "Fluent. Confident. Invented.\nThis failure is called HALLUCINATION.", 1.25, 4.48, 4.75, 0.6, 14, RED, True)

        shape_box(s, 6.85, 1.5, 5.78, 4.05, PALE, True, LINE, 1)
        shape_box(s, 6.85, 1.5, 5.78, 0.52, TEAL, True, TEAL, 0)
        text_box(s, "WITH RAG", 6.85, 1.62, 5.78, 0.3, 13, WHITE, True, PP_ALIGN.CENTER)
        text_box(s, "The workflow silently rewrites the question:", 7.2, 2.25, 5.1, 0.35, 15, GREY)
        shape_box(s, 7.2, 2.68, 5.08, 1.72, WHITE, True, TEAL, 1)
        text_box(s, "“Using ONLY these passages from the handbook:\n\n  <retrieved chunk 1>\n  <retrieved chunk 2>\n\nanswer: what is the refund policy?\nIf they do not contain the answer, say you do not know.”", 7.4, 2.78, 4.7, 1.55, 12, INK)
        text_box(s, "The model stops being a SOURCE of facts\nand becomes a READER of facts you supply.", 7.2, 4.55, 5.08, 0.7, 15, TEAL, True)

        # ---- 2. The two phases, as a pipeline diagram ----
        s = base("RAG Has Two Phases", "RAG - ARCHITECTURE")
        text_box(s, "Everything below exists to answer ONE question: which passages do we paste in front of the user's question?", 0.75, 1.42, 11.9, 0.4, 16, VIOLET, True)

        shape_box(s, 0.75, 2.05, 11.85, 1.85, PALE, True, LINE, 1)
        circle(s, "1", 1.1, 2.25, 0.5, BLUE, WHITE, 13)
        text_box(s, "INDEXING", 1.75, 2.28, 1.7, 0.3, 14, BLUE, True)
        text_box(s, "runs ONCE, when a document is uploaded", 3.3, 2.31, 4.2, 0.28, 12, GREY)
        idx_steps = [("PDF", 1.15), ("Split into\nchunks", 3.15), ("Embed each\nchunk", 5.6), ("768-dim\nvectors", 8.05), ("Vector\nstore", 10.5)]
        for i, (label, x) in enumerate(idx_steps):
            shape_box(s, x, 2.95, 1.75, 0.72, WHITE, True, BLUE, 1)
            text_box(s, label, x, 3.02, 1.75, 0.6, 12, INK, True, PP_ALIGN.CENTER)
            if i < len(idx_steps) - 1:
                arrow(s, x + 1.87, 3.31, BLUE)

        shape_box(s, 0.75, 4.15, 11.85, 2.15, PALE, True, LINE, 1)
        circle(s, "2", 1.1, 4.35, 0.5, TEAL, WHITE, 13)
        text_box(s, "RETRIEVAL", 1.75, 4.38, 1.9, 0.3, 14, TEAL, True)
        text_box(s, "runs on EVERY question", 3.5, 4.41, 3.6, 0.28, 12, GREY)
        ret_steps = [("Question", 1.15), ("Embed the\nquestion", 3.15), ("Cosine\nsimilarity", 5.6), ("Top-k\nchunks", 8.05), ("Grounded\nanswer", 10.5)]
        for i, (label, x) in enumerate(ret_steps):
            shape_box(s, x, 5.05, 1.75, 0.72, WHITE, True, TEAL, 1)
            text_box(s, label, x, 5.12, 1.75, 0.6, 12, INK, True, PP_ALIGN.CENTER)
            if i < len(ret_steps) - 1:
                arrow(s, x + 1.87, 5.41, TEAL)
        text_box(s, "On the n8n canvas these are the two visible paths: the upload path ending at the vector store, and the chat path reading from it.", 1.15, 5.9, 11.1, 0.3, 13, GREY)

        # ---- 3. Tokenization ----
        s = base("Tokenization: How Text Becomes Numbers", "RAG - TOKENS")
        text_box(s, "Models do not read words. They read TOKENS - the sub-word units in the model's vocabulary.", 0.75, 1.42, 11.9, 0.35, 16, GREY)
        text_box(s, "“The salon is closed on Sundays.”", 0.78, 1.95, 6.0, 0.42, 19, INK, True)
        tokens = ["The", " salon", " is", " closed", " on", " Sund", "ays", "."]
        x = 0.78
        for i, tok in enumerate(tokens):
            w = max(0.72, 0.30 + 0.145 * len(tok))
            color = AMBER if tok in (" Sund", "ays") else BLUE
            shape_box(s, x, 2.55, w, 0.62, WHITE, True, color, 1)
            text_box(s, tok.strip(), x, 2.64, w, 0.4, 14, color, True, PP_ALIGN.CENTER)
            x += w + 0.14
        text_box(s, "8 tokens  ·  note that “Sundays” split into TWO", 0.78, 3.3, 7.5, 0.3, 13, AMBER, True)

        shape_box(s, 0.75, 3.9, 11.85, 0.75, PALE, True, PALE, 0)
        text_box(s, "RULE OF THUMB (English):   1 token ≈ 4 characters ≈ 0.75 of a word   →   1,000 tokens ≈ 750 words ≈ 1.5 pages", 1.1, 3.95, 11.2, 0.65, 15, VIOLET, True, PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

        why = [
            ("CONTEXT WINDOW", "System prompt, retrieved chunks, memory and the question all compete for the SAME token budget.", BLUE),
            ("COST + LATENCY", "Double the retrieved text and you roughly double the prompt cost of every call.", AMBER),
            ("CHUNK SIZE", "Chunks are measured in tokens. This is the number you tune in Lab 2.4.", TEAL),
        ]
        for i, (label, note, color) in enumerate(why):
            x = 0.75 + i * 4.05
            shape_box(s, x, 4.95, 3.75, 1.5, WHITE, True, color, 1)
            shape_box(s, x, 4.95, 3.75, 0.42, color, True, color, 0)
            text_box(s, label, x, 5.03, 3.75, 0.28, 12, WHITE, True, PP_ALIGN.CENTER)
            text_box(s, note, x + 0.18, 5.5, 3.4, 0.85, 13, GREY)

        # ---- 4. Embeddings ----
        s = base("Embeddings: Meaning as Coordinates", "RAG - EMBEDDINGS")
        text_box(s, "An embedding turns ANY text into a fixed-length vector. Here: nomic-embed-text (Ollama) → 768 numbers.", 0.75, 1.4, 11.9, 0.32, 15, GREY)
        shape_box(s, 0.75, 1.95, 5.4, 0.72, WHITE, True, BLUE, 1)
        text_box(s, "“How do I reset my password?”", 0.95, 2.06, 5.0, 0.45, 15, INK, True, valign=MSO_ANCHOR.MIDDLE)
        arrow(s, 6.3, 2.28, BLUE)
        shape_box(s, 7.0, 1.95, 5.6, 0.72, PALE, True, BLUE, 1)
        text_box(s, "[ 0.021, -0.118, 0.334, … ]   768 numbers", 7.2, 2.06, 5.2, 0.45, 14, BLUE, True, valign=MSO_ANCHOR.MIDDLE)

        text_box(s, "Texts with similar MEANING land close together - even with zero words in common.", 0.75, 2.85, 11.9, 0.32, 16, VIOLET, True)
        pairs = [
            ("“reset my password”  vs  “I forgot my login credentials”", "0.89", "NEAR - no shared words, same meaning", TEAL),
            ("“reset my password”  vs  “what are your opening hours?”", "0.11", "FAR - unrelated", RED),
        ]
        for i, (pair, score, note, color) in enumerate(pairs):
            y = 3.35 + i * 0.92
            shape_box(s, 0.75, y, 8.3, 0.72, WHITE, True, LINE, 1)
            text_box(s, pair, 0.95, y + 0.1, 8.0, 0.5, 14, INK, False, valign=MSO_ANCHOR.MIDDLE)
            shape_box(s, 9.25, y, 1.3, 0.72, color, True, color, 0)
            text_box(s, score, 9.25, y + 0.16, 1.3, 0.4, 17, WHITE, True, PP_ALIGN.CENTER)
            text_box(s, note, 10.75, y + 0.1, 2.0, 0.55, 11, color, True, valign=MSO_ANCHOR.MIDDLE)
        text_box(s, "Cosine similarity:  1.0 = same direction/meaning   ·   0.0 = unrelated   ·   -1.0 = opposite. Retrieval returns the top-k nearest chunks (k is usually 3-5).", 0.75, 5.2, 11.9, 0.32, 13, GREY)

        shape_box(s, 0.75, 5.68, 11.85, 0.8, PALE, True, PALE, 0)
        text_box(s, "TWO RULES THAT CAUSE REAL BUGS:   Questions and documents MUST use the same embedding model · A vector store answers only “what is near this vector?” - not SQL", 1.05, 5.8, 11.3, 0.6, 14, RED, True, valign=MSO_ANCHOR.MIDDLE)

        # ---- 5. Chunking trade-off ----
        s = base("Chunking: The Trade-off You Must Measure", "RAG - CHUNKING")
        text_box(s, "You cannot embed a 40-page PDF as one vector - the detail would be averaged away. So documents are split into chunks (e.g. 800 characters) with an overlap (e.g. 100) so a sentence cut in half still appears whole somewhere.", 0.75, 1.42, 11.9, 0.6, 15, GREY)

        cols = [
            ("CHUNKS TOO SMALL", RED, [
                "Retrieved passage lacks context",
                "The model sees half an answer",
                "High precision, low recall",
                "“The document does not say” - when it does",
            ]),
            ("CHUNKS TOO LARGE", AMBER, [
                "Answer arrives with 3 irrelevant sections",
                "The model gets distracted, cites the wrong part",
                "High recall, low precision",
                "Tokens - and money - are wasted",
            ]),
        ]
        for i, (label, color, items) in enumerate(cols):
            x = 0.75 + i * 6.1
            shape_box(s, x, 2.2, 5.75, 2.9, PALE, True, LINE, 1)
            shape_box(s, x, 2.2, 5.75, 0.5, color, True, color, 0)
            text_box(s, label, x, 2.31, 5.75, 0.3, 13, WHITE, True, PP_ALIGN.CENTER)
            for j, item in enumerate(items):
                circle(s, "•", x + 0.35, 2.92 + j * 0.52, 0.26, color, WHITE, 10)
                text_box(s, item, x + 0.78, 2.9 + j * 0.52, 4.8, 0.42, 13, INK, False, valign=MSO_ANCHOR.MIDDLE)

        shape_box(s, 0.75, 5.3, 11.85, 1.15, PALE, True, PALE, 0)
        text_box(s, "There is no universally correct value.", 1.1, 5.45, 4.2, 0.3, 15, VIOLET, True)
        text_box(s, "That is why Lab 2.4 MEASURES it with golden questions instead of guessing. And when an answer is wrong: read the RETRIEVED CHUNKS first - if the right chunk never came back, no prompt rewrite will save it.", 1.1, 5.75, 11.2, 0.6, 14, INK, True)

    def ngrok_slides() -> None:
        # ---- Why a tunnel is needed at all ----
        s = base("Why localhost Stops Working", "TUNNEL - THE PROBLEM")
        text_box(s, "In Labs 4 and 5 something OUTSIDE your machine calls INTO n8n for the first time.", 0.75, 1.42, 11.9, 0.35, 17, GREY)

        shape_box(s, 0.72, 2.0, 5.9, 2.5, PALE, True, LINE, 1)
        shape_box(s, 0.72, 2.0, 5.9, 0.5, TEAL, True, TEAL, 0)
        text_box(s, "WORKS - your own browser", 0.72, 2.11, 5.9, 0.3, 13, WHITE, True, PP_ALIGN.CENTER)
        text_box(s, "Labs 1-3  ·  Book by Voice  ·  curl", 1.05, 2.72, 5.3, 0.35, 15, INK, True)
        text_box(s, "Your browser and n8n are on the SAME machine,\nso http://localhost:5678 resolves correctly.", 1.05, 3.15, 5.3, 0.9, 14, GREY)
        text_box(s, "No tunnel needed. Do not change this.", 1.05, 4.02, 5.3, 0.3, 14, TEAL, True)

        shape_box(s, 6.95, 2.0, 5.68, 2.5, PALE, True, LINE, 1)
        shape_box(s, 6.95, 2.0, 5.68, 0.5, RED, True, RED, 0)
        text_box(s, "FAILS - ElevenLabs / Vapi servers", 6.95, 2.11, 5.68, 0.3, 13, WHITE, True, PP_ALIGN.CENTER)
        text_box(s, "check_availability · book_appointment · Custom LLM", 7.28, 2.7, 5.1, 0.32, 12, INK, True)
        text_box(s, "To a datacenter in another country, “localhost”\nmeans ITS OWN machine. The request never\nleaves their building.", 7.28, 3.1, 5.1, 0.95, 14, GREY)
        text_box(s, "Tunnel required.", 7.28, 4.05, 5.1, 0.3, 14, RED, True)

        shape_box(s, 0.72, 4.8, 11.9, 1.5, PALE, True, PALE, 0)
        text_box(s, "THE SYMPTOM YOU WILL ACTUALLY SEE", 1.05, 4.95, 6.0, 0.3, 12, VIOLET, True)
        text_box(s, "Nina says “let me check that for you” … then silence. Nothing appears in the n8n executions list, because the\nrequest never arrived. That is not a prompt bug. It is a networking bug.", 1.05, 5.32, 11.2, 0.85, 15, INK, True)

        # ---- The five commands ----
        s = base("Expose n8n With ngrok", "TUNNEL - SETUP")
        steps = [
            ("1", "Install", "brew install ngrok\n(Windows: ngrok.com/download)", BLUE),
            ("2", "Authtoken", "Free signup, then:\nngrok config add-authtoken <TOKEN>", TEAL),
            ("3", "Start it", "ngrok http 5678\nLeave this window OPEN all session", VIOLET),
            ("4", "Read the URL", "Forwarding  https://<id>.ngrok-free.app\n-> http://localhost:5678", AMBER),
            ("5", "Inspect", "http://127.0.0.1:4040   <- bookmark this\nStatus tab = your public URL.  Inspect tab = every request as it arrives.", RED),
        ]
        for i, (num, label, body_text, color) in enumerate(steps):
            y = 1.48 + i * 0.98
            circle(s, num, 0.78, y + 0.08, 0.56, color, WHITE, 14)
            shape_box(s, 1.6, y, 11.0, 0.86, WHITE, True, LINE, 1)
            text_box(s, label, 1.85, y + 0.06, 2.0, 0.35, 15, INK, True)
            text_box(s, body_text, 4.0, y + 0.05, 8.4, 0.76, 12, GREY)
        text_box(s, "Prove it with curl BEFORE touching the voice platform - HTTP is far easier to debug than audio.", 0.78, 6.52, 11.8, 0.3, 14, VIOLET, True)

        # ---- Building the URL ----
        s = base("Building the URL You Paste Into the Vendor", "TUNNEL - THE URL")
        text_box(s, "You are gluing two halves together: ngrok supplies the address, the n8n Webhook node supplies the path.", 0.75, 1.42, 11.9, 0.35, 16, GREY)

        shape_box(s, 0.75, 2.0, 6.4, 0.85, PALE, True, BLUE, 1)
        text_box(s, "https://3af5-…ngrok-free.app", 0.95, 2.18, 6.0, 0.45, 16, BLUE, True, valign=MSO_ANCHOR.MIDDLE)
        text_box(s, "from ngrok", 0.95, 2.88, 6.0, 0.28, 12, GREY)
        shape_box(s, 7.45, 2.0, 5.2, 0.85, PALE, True, TEAL, 1)
        text_box(s, "/webhook/check-availability", 7.65, 2.18, 4.8, 0.45, 16, TEAL, True, valign=MSO_ANCHOR.MIDDLE)
        text_box(s, "from the n8n Webhook node", 7.65, 2.88, 4.8, 0.28, 12, GREY)

        shape_box(s, 0.75, 3.4, 11.85, 0.72, PALE, True, PALE, 0)
        text_box(s, "THE RULE:  take the node's PRODUCTION URL and swap http://localhost:5678 for your ngrok address.", 1.05, 3.5, 11.3, 0.52, 16, VIOLET, True, valign=MSO_ANCHOR.MIDDLE)

        traps = [
            ("The free URL changes on every restart", "Re-paste it and PUBLISH the agent again", AMBER),
            ("Test URL instead of Production URL", "Use /webhook/ - never /webhook-test/", RED),
            ("Workflow not published", "404 'not registered' - publish, then retest", BLUE),
            ("Setting WEBHOOK_URL on the container", "Do NOT. It breaks localhost for Labs 1-3", VIOLET),
        ]
        for i, (trap, fix, color) in enumerate(traps):
            y = 4.4 + i * 0.72
            circle(s, "!", 0.85, y, 0.42, color, WHITE, 12)
            text_box(s, trap, 1.5, y - 0.02, 5.4, 0.45, 14, INK, True, valign=MSO_ANCHOR.MIDDLE)
            text_box(s, fix, 7.1, y - 0.02, 5.5, 0.45, 14, GREY, False, valign=MSO_ANCHOR.MIDDLE)

    def video_slides() -> None:
        # ---- The pipeline, as a flow diagram ----
        s = base("How the Avatar Video Pipeline Works", "VIDEO - PIPELINE")
        text_box(s, "Every video lab is this same pipeline. Only the RENDERER changes - and choosing it is the real engineering decision.", 0.75, 1.4, 11.9, 0.35, 16, GREY)

        stages = [
            ("YOU", "type a TOPIC\n(raw facts)", BLUE),
            ("n8n", "Webhook receives it", TEAL),
            ("gemma4", "turns facts into\nSPOKEN COPY", VIOLET),
            ("RENDERER", "HeyGen / Wav2Lip /\nMuseTalk", AMBER),
            ("PAGE", "the video plays", RED),
        ]
        for i, (label, note, color) in enumerate(stages):
            x = 0.72 + i * 2.5
            shape_box(s, x, 2.05, 2.15, 1.5, WHITE, True, color, 1)
            shape_box(s, x, 2.05, 2.15, 0.45, color, True, color, 0)
            text_box(s, label, x, 2.14, 2.15, 0.3, 13, WHITE, True, PP_ALIGN.CENTER)
            text_box(s, note, x + 0.12, 2.62, 1.9, 0.85, 12, INK, False, PP_ALIGN.CENTER)
            if i < len(stages) - 1:
                arrow(s, x + 2.2, 2.78, LINE)

        shape_box(s, 0.72, 3.85, 11.9, 1.35, PALE, True, PALE, 0)
        text_box(s, "THE TWO TEXTS ARE NOT THE SAME TEXT", 1.05, 3.98, 6.0, 0.3, 12, RED, True)
        text_box(s, "You type a TOPIC: “Belgium beats USA 4-1.”   ->   gemma4 writes the SCRIPT: “Welcome back to Global Sports News!…”\nThe renderer speaks EVERY character - feed it raw facts and the avatar reads out a list.", 1.05, 4.32, 11.3, 0.78, 13, INK, True)

        shape_box(s, 0.72, 5.45, 5.85, 1.35, WHITE, True, LINE, 1)
        text_box(s, "Lab 7 - HeyGen (cloud)", 0.95, 5.55, 5.4, 0.3, 13, AMBER, True)
        text_box(s, "n8n returns a video_id at once; the PAGE POLLS for status.\nA 1-3 minute HTTP request is fragile.", 0.95, 5.88, 5.4, 0.75, 13, GREY)
        shape_box(s, 6.75, 5.45, 5.87, 1.35, WHITE, True, LINE, 1)
        text_box(s, "Lab 7-os - local render", 6.98, 5.55, 5.4, 0.3, 13, TEAL, True)
        text_box(s, "n8n returns the finished video_url in ONE response.\n~16 s, so the caller can simply wait.", 6.98, 5.88, 5.4, 0.75, 13, GREY)

        # ---- Three engines, side by side ----
        s = base("Three Lip-Sync Engines, One Photo", "VIDEO - THE TRADE-OFF")
        text_box(s, "Lab 6: MuseTalk vs HeyGen on one script and portrait. Lab 7-os adds Wav2Lip.", 0.75, 1.4, 11.9, 0.35, 16, GREY)

        cols = [
            ("WAV2LIP", TEAL, [
                ("Where", "your machine"),
                ("Cost", "free"),
                ("Speed", "~16 s   FASTEST"),
                ("Mouth", "96x96 - soft at 1080p"),
                ("Head", "frozen"),
                ("Needs", "ffmpeg + checkpoint"),
            ]),
            ("MUSETALK", VIOLET, [
                ("Where", "your machine"),
                ("Cost", "free"),
                ("Speed", "~75 s  (7x slower)"),
                ("Mouth", "real pixels  BEST"),
                ("Head", "frozen"),
                ("Needs", "GPU + 3.5 GB weights"),
            ]),
            ("HEYGEN", AMBER, [
                ("Where", "the cloud"),
                ("Cost", "CREDITS"),
                ("Speed", "~40 s - 3 min"),
                ("Mouth", "photoreal"),
                ("Head", "MOVES + blinks"),
                ("Needs", "API key + credits"),
            ]),
        ]
        for i, (name, color, rows) in enumerate(cols):
            x = 0.72 + i * 4.05
            shape_box(s, x, 1.95, 3.8, 4.15, PALE, True, LINE, 1)
            shape_box(s, x, 1.95, 3.8, 0.5, color, True, color, 0)
            text_box(s, name, x, 2.06, 3.8, 0.3, 14, WHITE, True, PP_ALIGN.CENTER)
            for j, (k, v) in enumerate(rows):
                y = 2.6 + j * 0.56
                text_box(s, k, x + 0.2, y, 1.15, 0.3, 11, GREY, False)
                text_box(s, v, x + 1.35, y, 2.35, 0.35, 12, INK, True)

        shape_box(s, 0.72, 6.22, 11.9, 0.82, PALE, True, PALE, 0)
        text_box(s, "Wav2Lip is FASTEST · MuseTalk looks BEST · only HeyGen MOVES THE HEAD", 1.05, 6.3, 11.3, 0.3, 14, VIOLET, True)
        text_box(s, "The question is not “which is best” - it is whether you would upload a customer's face to a vendor to gain a moving head.", 1.05, 6.63, 11.3, 0.3, 13, INK, False)

    def voice_arch_slide() -> None:
        """Topic 2 opener: the ElevenLabs-vs-Vapi split IS the lesson."""
        s = base("Two Vendors, Two Architectures", "VOICE - THE CONTRAST IS THE LESSON")
        cols = [
            ("ELEVENLABS - LAB 4 (GG HAIR SALON)", BLUE, [
                ("Who runs the model", "ElevenLabs"),
                ("What n8n does", "mints a signed URL, serves the tools"),
                ("Call path", "browser -> n8n -> signed URL -> WebSocket"),
                ("Browser receives", "a short-lived SIGNED URL"),
                ("Browser never sees", "your xi-api-key"),
            ]),
            ("VAPI - LAB 5 (MEDIREFILL)", TEAL, [
                ("Who runs the model", "YOUR n8n workflow (Custom LLM)"),
                ("What n8n does", "IS the brain - and the guardrail"),
                ("Call path", "browser -> Vapi, then Vapi -> your n8n"),
                ("Browser receives", "the Vapi PUBLIC key only"),
                ("Browser never sees", "your Vapi private key"),
            ]),
        ]
        for i, (name, color, rows) in enumerate(cols):
            x = 0.72 + i * 6.1
            shape_box(s, x, 1.75, 5.85, 4.35, PALE, True, LINE, 1)
            shape_box(s, x, 1.75, 5.85, 0.5, color, True, color, 0)
            text_box(s, name, x, 1.86, 5.85, 0.3, 13, WHITE, True, PP_ALIGN.CENTER)
            for j, (k, v) in enumerate(rows):
                y = 2.5 + j * 0.7
                text_box(s, k, x + 0.25, y, 2.1, 0.6, 11, GREY, False)
                text_box(s, v, x + 2.45, y, 3.25, 0.6, 12, INK, True)
        shape_box(s, 0.72, 6.22, 11.9, 0.68, PALE, True, PALE, 0)
        text_box(s, "Either way, the browser only ever holds a short-lived, single-purpose token - never a vendor key.",
                 1.05, 6.32, 11.3, 0.5, 14, VIOLET, True, valign=MSO_ANCHOR.MIDDLE)

    for idx, (topic_title, topic_desc) in enumerate(TOPICS, 1):
        topic_slide(idx, topic_title, topic_desc)
        if idx == 1:
            rag_concept_slides()
        if idx == 2:
            voice_arch_slide()
            ngrok_slides()
        if idx == 3:
            video_slides()
        for lab in [l for l in LABS if l.topic == idx]:
            # Record the slide number BEFORE the lab's first slide is added, so the
            # Lesson Plan cites the deck as it actually is. Any reorder of the deck
            # updates the LP automatically — they cannot drift apart.
            SLIDE_MAP[lab.lab_no] = len(prs.slides._sldIdLst) + 1
            lab_target(lab)
            concept_slide(lab)
            step_slide(lab)
            # The screenshots belong WITH the lab, not in a gallery at the end of
            # the deck - a trainer walking through Lab 4 must see the canvas here.
            for key in lab_shots(lab):
                screenshot_slide(key, f"LAB {lab.lab_no} - {lab.title.upper()}")
            verify_slide(lab)
            fixes_slide(lab)

    s = base("Four Businesses, Four Agents", "YOUR CASE STUDY IS BUILT FROM THE LABS")
    businesses = [
        ("COOK & BAKE ACADEMY", "Website course advisor, grounded ONLY in the brochures - honest about what it does not know.", "Lab 3", BLUE),
        ("GG HAIR SALON", "Voice receptionist who checks the REAL Google Calendar and books into it.", "Lab 4", VIOLET),
        ("MEDIREFILL", "Refill voice FAQ that never gives medical advice - refusal + pharmacist callback.", "Lab 5", TEAL),
        ("GG NEWS STUDIO", "Presenter video from the day's facts - and a cloud vs open-source recommendation.", "Labs 7 / 7-os", AMBER),
    ]
    for i, (name, note, lab_ref, color) in enumerate(businesses):
        x = 0.72 + (i % 2) * 6.1
        y = 1.6 + (i // 2) * 2.15
        shape_box(s, x, y, 5.85, 1.92, PALE, True, LINE, 1)
        shape_box(s, x, y, 5.85, 0.46, color, True, color, 0)
        text_box(s, name, x, y + 0.09, 4.3, 0.3, 13, WHITE, True, PP_ALIGN.CENTER)
        shape_box(s, x + 4.35, y + 0.06, 1.35, 0.34, WHITE, True, WHITE, 0)
        text_box(s, lab_ref, x + 4.35, y + 0.11, 1.35, 0.26, 11, color, True, PP_ALIGN.CENTER)
        text_box(s, note, x + 0.3, y + 0.66, 5.3, 1.15, 14, INK, False, valign=MSO_ANCHOR.MIDDLE)
    shape_box(s, 0.72, 6.05, 11.9, 0.85, PALE, True, PALE, 0)
    text_box(s, "The Case Study asks you to rebuild and defend exactly these - keep your workflow exports, traces and transcripts.",
             1.05, 6.19, 11.3, 0.6, 15, VIOLET, True, valign=MSO_ANCHOR.MIDDLE)

    s = base("Assessment Rubric", "EVIDENCE STANDARD")
    rubric = [
        ("Problem definition", "User, trigger, output and non-goals", BLUE),
        ("Workflow correctness", "Expected paths proven by execution", TEAL),
        ("AI quality", "Prompts, retrieval and media evaluated", VIOLET),
        ("Guardrails", "Secrets, claims and publishing controlled", AMBER),
        ("Documentation", "Another operator can run and recover", RED),
    ]
    for i, (label, note, color) in enumerate(rubric):
        y = 1.5 + i * 1.0
        circle(s, str(i + 1), 0.78, y, 0.54, color, WHITE, 13)
        text_box(s, label, 1.58, y - 0.02, 2.35, 0.32, 16, INK, True)
        text_box(s, note, 4.0, y - 0.02, 4.95, 0.38, 14, GREY)
        shape_box(s, 9.28, y + 0.06, 2.8, 0.18, PALE, True, PALE, 0)
        shape_box(s, 9.28, y + 0.06, 0.56 * (i + 1), 0.18, color, True, color, 0)
    text_box(s, "Score the evidence, explain one improvement, then rerun the test.", 1.58, 6.45, 10.2, 0.32, 16, VIOLET, True)

    # ── closing block ────────────────────────────────────────────────────────────
    # Briefing -> Assessment -> Assessment Flow -> TRAQOM -> Thank You.
    # This course has NO practice exam - do not add a Practice Exam slide.
    s = base("Briefing for Assessment", "BEFORE YOU ARE ASSESSED")
    text_box(s, "Read this before the papers are handed out. You are entitled to know exactly how you will be judged.",
             0.75, 1.42, 11.9, 0.4, 17, GREY)
    briefing = [
        ("What you sit", "A Written Assessment (short answer) and a Case Study.", BLUE),
        ("How long", "60 minutes each. Open book. Individually attempted.", TEAL),
        ("How it is marked", "Competent / Not Yet Competent against the assessment criteria.", VIOLET),
        ("If you are NYC", "You are re-assessed once, on the failed instrument only.", AMBER),
        ("Appeals", "Tell the trainer on the day; the appeal route is in your Learner Guide.", RED),
    ]
    for i, (label, note, color) in enumerate(briefing):
        y = 2.05 + i * 0.95
        circle(s, str(i + 1), 0.85, y, 0.58, color, WHITE, 14)
        text_box(s, label, 1.75, y - 0.02, 3.1, 0.4, 17, INK, True, valign=MSO_ANCHOR.MIDDLE)
        text_box(s, note, 5.0, y - 0.02, 7.6, 0.5, 15, GREY, False, valign=MSO_ANCHOR.MIDDLE)

    s = base("Assessment", "HOW YOU ARE ASSESSED")
    for i, (name, dur, tests, color) in enumerate([
        ("Written Assessment (WA)", "60 min", "KNOWLEDGE - the concepts behind the labs: RAG, grounding, webhooks, lip-sync engines, guardrails.", BLUE),
        ("Case Study (CS)", "60 min", "APPLICATION - judgement on a workplace scenario built from the labs you ran.", VIOLET),
    ]):
        x = 0.75 + i * 6.1
        shape_box(s, x, 1.6, 5.85, 4.2, PALE, True, LINE, 1)
        shape_box(s, x, 1.6, 5.85, 0.6, color, True, color, 0)
        text_box(s, name, x, 1.73, 5.85, 0.35, 16, WHITE, True, PP_ALIGN.CENTER)
        circle(s, dur, x + 2.35, 2.4, 1.15, color, WHITE, 15)
        text_box(s, tests, x + 0.35, 3.9, 5.15, 1.6, 15, INK, False, PP_ALIGN.CENTER)
    shape_box(s, 0.75, 6.05, 11.85, 0.95, PALE, True, PALE, 0)
    text_box(s, "Open book · individually attempted · marked Competent / Not Yet Competent · one re-assessment on the failed instrument",
             1.05, 6.2, 11.3, 0.6, 15, VIOLET, True, PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    s = base("Assessment Flow", "FROM PAPER TO SIGN-OFF")
    # The full path, not just the papers: attendance is taken for the assessment session,
    # the scripts are submitted on the LMS, and the day ends on the Assessment Summary
    # Record. Leave those out and the diagram describes an exam, not a WSQ assessment.
    flow = [
        ("BRIEFING", "criteria explained", BLUE),
        ("ATTENDANCE", "TRAQOM, digital", TEAL),
        ("WA (SAQ)", "60 min · knowledge", BLUE),
        ("CASE STUDY", "60 min · application", VIOLET),
        ("SUBMIT", "upload on the LMS", TEAL),
        ("MARKING", "C / NYC per criterion", AMBER),
        ("SIGN-OFF", "Assessment Summary Record", RED),
    ]
    for i, (label, note, color) in enumerate(flow):
        x = 0.62 + i * 1.82
        shape_box(s, x, 2.3, 1.62, 1.75, WHITE, True, color, 1)
        shape_box(s, x, 2.3, 1.62, 0.44, color, True, color, 0)
        text_box(s, label, x, 2.37, 1.62, 0.3, 10, WHITE, True, PP_ALIGN.CENTER)
        text_box(s, note, x + 0.08, 2.88, 1.46, 1.0, 10, INK, False, PP_ALIGN.CENTER)
        if i < len(flow) - 1:
            arrow(s, x + 1.66, 3.14, LINE)
    shape_box(s, 0.62, 4.5, 12.0, 1.5, PALE, True, PALE, 0)
    text_box(s, "NOT YET COMPETENT?", 0.95, 4.65, 4.0, 0.3, 12, RED, True)
    text_box(s, "You are re-assessed ONCE, on the failed instrument only. The trainer records the gap, you close it, and you sit that paper again.\nA re-assessment is a normal part of competency-based training - it is not a failure of the course.",
             0.95, 5.0, 11.4, 0.85, 15, INK, True)

    traqom_slide("at the END")

    s = base("Thank You", "COURSE CLOSE")
    shape_box(s, 0.75, 1.6, 11.82, 3.5, PALE, True, PALE, 0)
    text_box(s, "A polished output is only the beginning.", 1.2, 2.05, 10.9, 0.55, 27, INK, True, PP_ALIGN.CENTER)
    text_box(s, "Professional agentic systems are observable, testable, guarded, recoverable, and documented.", 1.5, 2.95, 10.3, 0.88, 21, GREY, False, PP_ALIGN.CENTER)
    for i, (label, color) in enumerate([("TRACE", BLUE), ("TEST", VIOLET), ("REVIEW", TEAL), ("RECOVER", AMBER), ("DOCUMENT", RED)]):
        x = 1.25 + i * 2.18
        shape_box(s, x, 4.42, 1.82, 0.52, color, True, color, 0)
        text_box(s, label, x, 4.52, 1.82, 0.24, 11, WHITE, True, PP_ALIGN.CENTER)
    text_box(s, "www.tertiarycourses.com.sg", 4.4, 5.82, 4.5, 0.38, 16, BLUE, True, PP_ALIGN.CENTER)

    # Add footers last so later content cannot mask them in PowerPoint or LibreOffice.
    for slide_no, slide in enumerate(prs.slides, 1):
        add_footer(slide, slide_no)

    prs.save(path)
    return SLIDE_MAP


def toc_entries_from_pdf(md: str, pdf: Path):
    """(level, heading, page) for every heading, with the page it ACTUALLY landed on.

    Read out of the rendered PDF rather than guessed, so the contents page agrees with the
    document a learner is holding. Returns [] if the PDF or PyMuPDF is unavailable - the
    caller then keeps the plain field, which Word will still populate on open.
    """
    try:
        import fitz
    except Exception:
        return []
    if not pdf.exists():
        return []

    heads = []
    for line in md.splitlines():
        for hashes, level in (("### ", 3), ("## ", 2), ("# ", 1)):
            if line.startswith(hashes):
                text = plain(line[len(hashes):]).strip()
                if text:
                    heads.append((level, text))
                break

    doc = fitz.open(pdf)
    pages = [p.get_text() for p in doc]
    doc.close()

    # Skip the contents pages themselves. Once the TOC carries a real result, every heading
    # ALSO appears on the TOC page - so a naive forward scan finds "Course profile" on the
    # contents page and reports page 4 for the whole document. It even converges there: the
    # wrong answer reproduces itself on the next pass.
    body_start = 0
    for i, text in enumerate(pages):
        if "TABLE OF CONTENTS" in text.upper():
            body_start = i + 1
            while body_start < len(pages) and "......" in pages[body_start]:
                body_start += 1
            break

    entries, start = [], body_start
    for level, text in heads:
        # Headings are unique and in order, so scan forward only: a lab title that repeats
        # in a cross-reference must not drag the entry back to an earlier page.
        needle = text.lower()
        for i in range(start, len(pages)):
            if needle in " ".join(pages[i].lower().split()):
                entries.append((level, text, i + 1))
                start = i
                break
    return entries


def to_pdf(path: Path, quiet: bool = False) -> None:
    """LibreOffice renders the PDF. A missing PDF is a QA failure, so say so loudly."""
    import subprocess

    try:
        subprocess.run(
            ["soffice", "--headless", "--convert-to", "pdf", "--outdir", str(path.parent), str(path)],
            check=True, capture_output=True, timeout=300,
        )
        pdf = path.with_suffix(".pdf")
        if quiet:
            return
        print(f"  PDF: {pdf.name} ({pdf.stat().st_size // 1024} KB)" if pdf.exists()
              else f"  PDF FAILED for {path.name}")
    except Exception as exc:
        print(f"  PDF skipped for {path.name}: {exc}")


def archive_old(keep: set[str]) -> None:
    """One live version in courseware/, everything superseded in courseware/archive/."""
    ARCHIVE.mkdir(parents=True, exist_ok=True)
    moved = 0
    for f in COURSEWARE.iterdir():
        if f.is_dir() or f.name in keep or f.name.startswith("~$"):
            continue
        if f.suffix.lower() in {".pptx", ".docx", ".pdf", ".md"}:
            f.rename(ARCHIVE / f.name)
            moved += 1
    if moved:
        print(f"  archived {moved} superseded file(s) -> courseware/archive/")


def main() -> None:
    COURSEWARE.mkdir(exist_ok=True)

    # The labs are REAL, runnable folders (labs_local_n8n/lab0 … lab10) - this
    # script documents them, it does not generate them. Refuse to build courseware
    # that references a lab folder which does not exist.
    missing_labs = [
        lab.lab_no for lab in LABS
        if not (LABS_DIR / f"lab{lab.lab_no.replace('-os', '-opensource')}").is_dir()
    ]
    assert not missing_labs, f"labs_local_n8n is missing lab folder(s): {missing_labs}"

    ppt = COURSEWARE / f"{PPT_STEM}.pptx"
    lg_docx = COURSEWARE / f"LG-{COURSE_TITLE.replace(' ', '-')}.docx"
    lp_docx = COURSEWARE / f"LP-{COURSE_TITLE.replace(' ', '-')}.docx"
    lg_md = COURSEWARE / f"LG-{COURSE_TITLE.replace(' ', '-')}.md"
    lp_md = COURSEWARE / f"LP-{COURSE_TITLE.replace(' ', '-')}.md"

    # The deck goes FIRST: it reports the slide number each lab starts on, and the
    # Lesson Plan must cite the deck as it actually is.
    slide_map = write_pptx(ppt)
    print(f"PPTX: {ppt.name}")

    # The repo root keeps its hand-written README.md and the two build-specific
    # guides (LEARNER_GUIDE_LOCAL.md / LEARNER_GUIDE_CLOUD.md) - do not overwrite
    # them. The generated LG/LP live in courseware/ only.
    lg_root = learner_guide("courseware/screenshots/")
    lg_courseware = learner_guide("screenshots/")
    lg_md.write_text(lg_courseware, encoding="utf-8")

    lp = lesson_plan(slide_map)
    lp_md.write_text(lp, encoding="utf-8")

    # The contents page cannot know its own page numbers until the document has been laid
    # out - and writing it CHANGES that layout, because a 60-entry TOC is itself pages long.
    # So render, read the real page of every heading, rewrite, and repeat until the numbers
    # stop moving. Two passes are not enough: they leave the TOC pointing one page short.
    for md, docx_path, kind in ((lg_root, lg_docx, "Learner Guide"), (lp, lp_docx, "LESSON PLAN")):
        entries = None
        for attempt in range(4):
            write_docx(md, docx_path, kind=kind, toc_entries=entries)
            to_pdf(docx_path, quiet=True)
            fresh = toc_entries_from_pdf(md, docx_path.with_suffix(".pdf"))
            if fresh == entries:
                break
            entries = fresh
        write_docx(md, docx_path, kind=kind, toc_entries=entries)
        print(f"DOCX: {docx_path.name}  (TOC: {len(entries or [])} entries, settled in {attempt + 1})")

    for f in (ppt, lg_docx, lp_docx):
        to_pdf(f)

    archive_old(keep={
        ppt.name, lg_docx.name, lp_docx.name, lg_md.name, lp_md.name,
        ppt.with_suffix(".pdf").name, lg_docx.with_suffix(".pdf").name, lp_docx.with_suffix(".pdf").name,
    })

    missing = [f for _, (f, _c) in SCREENSHOTS.items() if not (SHOTS_DIR / f).exists()]
    if missing:
        print(f"\nScreenshots still to capture ({len(missing)}): {', '.join(sorted(missing))}")


if __name__ == "__main__":
    main()
